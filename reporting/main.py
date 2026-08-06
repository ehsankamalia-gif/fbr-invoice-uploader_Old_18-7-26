import os
import json
import time
import threading
import logging
import smtplib
import secrets
from email.message import EmailMessage
from datetime import datetime, timedelta
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from fastapi import Body, Depends, FastAPI, Header, HTTPException, Query, Request
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, JSONResponse, RedirectResponse, StreamingResponse
from sqlalchemy import asc, desc, func
from sqlalchemy.orm import Session, joinedload

from app.db.session import SessionLocal
from app.db.models import (
    Customer,
    Invoice,
    InvoiceItem,
    Motorcycle,
    ProductModel,
    ReportTemplate,
    ReportSchedule,
    ReportRun,
    AuditLog,
    PrintTemplateLayout,
)
from app.services.invoice_service import invoice_service
from reporting.lookup_utils import format_cnic, validate_lookup_inputs
from reporting.invoice_detail_utils import invoice_to_detail_dict

try:
    from app.services.fastreport_bridge import (
        build_report as _fr_build_report,
        is_fastreports_available as _fr_available,
    )
except Exception:  # pragma: no cover - bridge import guard
    _fr_available = None

    def _fr_build_report(*_a, **_kw):
        from dataclasses import make_dataclass
        Br = make_dataclass("Br", [("ok", bool), ("output_path", "Optional[Path]"), ("error", "Optional[str]"), ("args_used", list)])
        return Br(ok=False, output_path=None, error="bridge_import_unavailable", args_used=[])

logger = logging.getLogger(__name__)


def get_db() -> Session:
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def _reporting_root_dir() -> Path:
    return Path(__file__).resolve().parent.parent / "exports" / "reports"


def _get_role(role_header: Optional[str]) -> str:
    role = (role_header or "").strip().lower()
    return role or "sales"


def _require_auth(
    api_key: Optional[str],
    role: str,
    required_roles: Optional[List[str]] = None,
) -> None:
    configured_key = (os.getenv("REPORTING_ACCESS_TOKEN") or "").strip()
    if configured_key and (api_key or "").strip() != configured_key:
        raise HTTPException(status_code=401, detail="Unauthorized")
    if required_roles and role not in [r.lower() for r in required_roles]:
        raise HTTPException(status_code=403, detail="Forbidden")


def _audit(
    db: Session,
    *,
    action: str,
    resource_type: str,
    resource_id: int,
    details: Dict[str, Any],
    request: Optional[Request] = None,
) -> None:
    ip = None
    if request:
        ip = request.client.host if request.client else None
    db.add(
        AuditLog(
            user_id=None,
            action=action,
            resource_type=resource_type,
            resource_id=resource_id,
            details=details,
            ip_address=ip,
        )
    )
    db.commit()


app = FastAPI(title="FBR Reporting Portal v2")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

_print_preview_lock = threading.Lock()
_print_previews: Dict[str, Dict[str, Any]] = {}


def _purge_expired_previews(now: float) -> None:
    with _print_preview_lock:
        expired = [k for k, v in _print_previews.items() if float(v.get("expires_at", 0.0) or 0.0) <= now]
        for k in expired:
            _print_previews.pop(k, None)


@app.post("/api/print-preview")
def api_create_print_preview(payload: Dict[str, Any] = Body(default={})) -> JSONResponse:
    html = ""
    title = ""
    if isinstance(payload, dict):
        html = str(payload.get("html") or "")
        title = str(payload.get("title") or "")
    if not html.strip():
        raise HTTPException(status_code=400, detail="html required")

    _purge_expired_previews(time.time())
    preview_id = secrets.token_urlsafe(12)
    expires_at = time.time() + 60 * 20
    with _print_preview_lock:
        _print_previews[preview_id] = {"html": html, "title": title, "expires_at": expires_at}
    return JSONResponse({"id": preview_id})


@app.get("/print-preview/{preview_id}", response_class=HTMLResponse)
def print_preview_page(preview_id: str) -> str:
    _purge_expired_previews(time.time())
    with _print_preview_lock:
        entry = _print_previews.get(preview_id)
    if not entry:
        return "<html><body><h3>Preview expired</h3></body></html>"
    return str(entry.get("html") or "")


@app.get("/api/print-layout/{template_name}")
def get_print_layout(
    template_name: str,
    db: Session = Depends(get_db),
) -> JSONResponse:
    key = (template_name or "").strip().lower()
    if not key:
        raise HTTPException(status_code=400, detail="template_name required")

    try:
        from app.db.session import engine as _engine
        logger.info(f"Loading print layout '{key}' using DB: {_engine.url}")
        row = db.query(PrintTemplateLayout).filter(PrintTemplateLayout.template_name == key).first()
        if not row:
            return JSONResponse({"template_name": key, "positions": {"version": 2, "updated_at": int(time.time() * 1000), "elements": {}}})
        return JSONResponse({"template_name": key, "positions": row.positions or {"version": 2, "updated_at": int(time.time() * 1000), "elements": {}}})
    except Exception as exc:
        logger.error(f"Failed to load print layout '{key}': {exc}", exc_info=True)
        return JSONResponse({"template_name": key, "positions": {"version": 2, "updated_at": int(time.time() * 1000), "elements": {}}, "error": "load_failed"})


@app.post("/api/print-layout/{template_name}")
def save_print_layout(
    template_name: str,
    payload: Dict[str, Any] = Body(default={}),
    db: Session = Depends(get_db),
) -> JSONResponse:
    key = (template_name or "").strip().lower()
    if not key:
        raise HTTPException(status_code=400, detail="template_name required")

    positions = payload.get("positions") if isinstance(payload, dict) else None
    if not isinstance(positions, dict):
        raise HTTPException(status_code=400, detail="payload.positions must be an object")
    if positions.get("version") != 2:
        raise HTTPException(status_code=400, detail="positions.version must be 2")
    elements = positions.get("elements")
    if elements is None:
        positions["elements"] = {}
    if positions.get("updated_at") is None:
        positions["updated_at"] = int(time.time() * 1000)

    try:
        from app.db.session import engine as _engine
        logger.info(f"Saving print layout '{key}' using DB: {_engine.url}")
        row = db.query(PrintTemplateLayout).filter(PrintTemplateLayout.template_name == key).first()
        if not row:
            row = PrintTemplateLayout(template_name=key, positions=positions)
            db.add(row)
        else:
            row.positions = positions
        db.commit()
        return JSONResponse({"ok": True, "template_name": key})
    except Exception as exc:
        db.rollback()
        logger.error(f"Failed to save print layout '{key}': {exc}", exc_info=True)
        raise HTTPException(status_code=500, detail="save_failed")

_retry_lock = threading.Lock()
_last_retry_at: Dict[str, float] = {}


DEFAULT_TEMPLATE = {
    "version": 1,
    "widgets": [
        {"type": "kpi", "metric": "total_invoices", "title": "Total Invoices"},
        {"type": "kpi", "metric": "total_amount", "title": "Total Amount"},
        {"type": "kpi", "metric": "avg_invoice_amount", "title": "Avg Invoice"},
        {"type": "chart", "metric": "daily_sales", "title": "Daily Sales"},
        {"type": "chart", "metric": "status_breakdown", "title": "Status Breakdown"},
        {"type": "table", "metric": "invoices", "title": "Invoices"},
    ],
}


def _parse_dates(from_date: Optional[str], to_date: Optional[str]) -> Tuple[Optional[datetime], Optional[datetime]]:
    start_dt = None
    end_dt = None
    if from_date:
        try:
            start_dt = datetime.strptime(from_date, "%Y-%m-%d")
        except ValueError:
            start_dt = None
    if to_date:
        try:
            end_dt = datetime.strptime(to_date, "%Y-%m-%d") + timedelta(days=1) - timedelta(seconds=1)
        except ValueError:
            end_dt = None
    return start_dt, end_dt


def _invoice_query(db: Session, start_dt: Optional[datetime], end_dt: Optional[datetime], status: str) -> Any:
    q = db.query(Invoice)
    if start_dt:
        q = q.filter(Invoice.datetime >= start_dt)
    if end_dt:
        q = q.filter(Invoice.datetime <= end_dt)
    st = (status or "ALL").upper()
    if st != "ALL":
        q = q.filter(Invoice.sync_status == st)
    return q


def _compute_metrics(db: Session, start_dt: Optional[datetime], end_dt: Optional[datetime], status: str) -> Dict[str, Any]:
    q = _invoice_query(db, start_dt, end_dt, status)
    total_invoices = q.count()
    total_amount = q.with_entities(func.coalesce(func.sum(Invoice.total_amount), 0.0)).scalar() or 0.0
    avg_invoice_amount = float(total_amount) / float(total_invoices) if total_invoices else 0.0

    daily_rows = (
        q.with_entities(func.date(Invoice.datetime), func.coalesce(func.sum(Invoice.total_amount), 0.0))
        .group_by(func.date(Invoice.datetime))
        .order_by(func.date(Invoice.datetime))
        .all()
    )
    daily_sales = [{"date": str(d), "amount": float(a or 0.0)} for d, a in daily_rows]

    status_rows = (
        q.with_entities(Invoice.sync_status, func.count(Invoice.id))
        .group_by(Invoice.sync_status)
        .order_by(Invoice.sync_status)
        .all()
    )
    status_breakdown = [{"status": str(s or ""), "count": int(c or 0)} for s, c in status_rows]

    invoices = (
        q.order_by(desc(Invoice.datetime))
        .limit(200)
        .all()
    )
    invoice_rows = []
    for inv in invoices:
        effective_usin = inv.fbr_invoice_number or inv.usin or ""
        invoice_rows.append(
            {
                "invoice_number": inv.invoice_number,
                "datetime": inv.datetime.isoformat(sep=" ") if inv.datetime else "",
                "pos_id": inv.pos_id or "",
                "payment_mode": inv.payment_mode or "",
                "total_amount": float(inv.total_amount or 0),
                "sync_status": inv.sync_status,
                "fbr_invoice_number": inv.fbr_invoice_number or "",
                "usin": effective_usin,
            }
        )

    return {
        "total_invoices": int(total_invoices),
        "total_amount": float(total_amount),
        "avg_invoice_amount": float(avg_invoice_amount),
        "daily_sales": daily_sales,
        "status_breakdown": status_breakdown,
        "invoices": invoice_rows,
    }


def _load_or_create_default_template(db: Session) -> ReportTemplate:
    tmpl = db.query(ReportTemplate).filter(ReportTemplate.name == "Default Dashboard").first()
    if tmpl:
        return tmpl
    tmpl = ReportTemplate(name="Default Dashboard", description="Default reporting dashboard", definition=DEFAULT_TEMPLATE, is_active=True)
    db.add(tmpl)
    db.commit()
    db.refresh(tmpl)
    return tmpl


def _ensure_template_has_widgets(db: Session, tmpl: Optional[ReportTemplate]) -> ReportTemplate:
    """Return the template with a guaranteed non-empty widget list.

    If the template is missing or its definition has no widgets (e.g. it was
    created before the New-template fix and never saved via Builder), fall
    back to DEFAULT_TEMPLATE widgets in-memory so the dashboard / exports
    render usable content instead of an empty page.
    """
    if tmpl is None:
        return _load_or_create_default_template(db)
    widgets = ((tmpl.definition or {}).get("widgets") or [])
    if len(widgets) == 0:
        if not tmpl.definition:
            tmpl.definition = {}
        tmpl.definition["widgets"] = list(DEFAULT_TEMPLATE["widgets"])
    return tmpl


def _render_dashboard_html() -> str:
    return """
    <!doctype html>
    <html lang="en">
    <head>
      <meta charset="utf-8"/>
      <meta name="viewport" content="width=device-width, initial-scale=1"/>
      <title>Reporting Portal</title>
      <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
      <script src="https://cdn.jsdelivr.net/npm/plotly.js-dist@2.32.0/plotly.min.js"></script>
      <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/js/bootstrap.bundle.min.js"></script>
      <style>
        .date-field { position: relative; }
        .date-popover {
          position: absolute;
          top: calc(100% + 6px);
          left: 0;
          z-index: 1056;
          width: 292px;
          background: #fff;
          border: 1px solid rgba(0,0,0,.15);
          border-radius: .5rem;
          box-shadow: 0 .5rem 1rem rgba(0,0,0,.15);
          padding: .5rem;
        }
        .date-popover .dp-head { display: flex; align-items: center; justify-content: space-between; gap: .25rem; margin-bottom: .25rem; }
        .date-popover .dp-title { font-weight: 600; }
        .date-popover .dp-btn { border: 1px solid rgba(0,0,0,.15); background: #f8f9fa; border-radius: .375rem; padding: .25rem .5rem; }
        .date-popover .dp-grid { display: grid; grid-template-columns: repeat(7, 1fr); gap: .15rem; }
        .date-popover .dp-dow { font-size: .75rem; color: #6c757d; text-align: center; padding: .25rem 0; }
        .date-popover .dp-day {
          border: 1px solid transparent;
          background: transparent;
          border-radius: .375rem;
          padding: .35rem 0;
          text-align: center;
          cursor: pointer;
        }
        .date-popover .dp-day:hover { background: #eef5ff; }
        .date-popover .dp-day[aria-selected="true"] { background: #0d6efd; color: #fff; }
        .date-popover .dp-day.dp-muted { color: #adb5bd; }
        .date-input.is-invalid { border-color: #dc3545; padding-right: calc(1.5em + .75rem); }
        .date-input.is-valid { border-color: #198754; padding-right: calc(1.5em + .75rem); }
      </style>
    </head>
    <body class="bg-light">
      <nav class="navbar navbar-expand-lg navbar-dark bg-primary">
        <div class="container-fluid">
          <span class="navbar-brand">Reporting Portal</span>
          <div class="d-flex gap-2">
            <a class="btn btn-outline-light btn-sm" href="/builder">FastReport Studio</a>
            <a class="btn btn-outline-light btn-sm" href="/schedules">Schedules</a>
            <a class="btn btn-outline-light btn-sm" href="/lookup">Lookup</a>
          </div>
        </div>
      </nav>
      <main class="container-fluid py-3">
        <div class="row g-3">
          <div class="col-12">
            <div class="card">
              <div class="card-body">
                <div class="row g-2 align-items-end">
                  <div class="col-12 col-md-3">
                    <label class="form-label">Template</label>
                    <select class="form-select" id="templateSelect"></select>
                  </div>
                  <div class="col-6 col-md-2">
                    <label class="form-label">From</label>
                    <div class="input-group date-field" data-date-field>
                      <input class="form-control date-input" type="text" id="fromDate" autocomplete="off" inputmode="text"
                             aria-label="From date" placeholder="YYYY-MM-DD or DD/MM/YYYY or MM/DD/YYYY"/>
                      <button class="btn btn-outline-secondary date-clear" type="button" aria-label="Clear date">×</button>
                      <button class="btn btn-outline-secondary date-open" type="button" aria-label="Open calendar">📅</button>
                    </div>
                    <div class="invalid-feedback d-block" id="fromDateError" style="display:none;"></div>
                  </div>
                  <div class="col-6 col-md-2">
                    <label class="form-label">To</label>
                    <div class="input-group date-field" data-date-field>
                      <input class="form-control date-input" type="text" id="toDate" autocomplete="off" inputmode="text"
                             aria-label="To date" placeholder="YYYY-MM-DD or DD/MM/YYYY or MM/DD/YYYY"/>
                      <button class="btn btn-outline-secondary date-clear" type="button" aria-label="Clear date">×</button>
                      <button class="btn btn-outline-secondary date-open" type="button" aria-label="Open calendar">📅</button>
                    </div>
                    <div class="invalid-feedback d-block" id="toDateError" style="display:none;"></div>
                  </div>
                  <div class="col-12 col-md-2">
                    <label class="form-label">Status</label>
                    <select class="form-select" id="status">
                      <option value="ALL">All</option>
                      <option value="PENDING">Pending</option>
                      <option value="SUCCESS">Success</option>
                      <option value="FAILED">Failed</option>
                    </select>
                  </div>
                  <div class="col-12 col-md-3 d-flex gap-2">
                    <button class="btn btn-primary flex-grow-1" id="applyBtn">Apply</button>
                    <button class="btn btn-outline-secondary" id="autoBtn">Auto</button>
                  </div>
                </div>
                <div class="mt-3 d-flex flex-wrap gap-2">
                  <button class="btn btn-outline-primary btn-sm" id="exportCsv">CSV</button>
                  <button class="btn btn-outline-primary btn-sm" id="exportXlsx">Excel</button>
                  <button class="btn btn-outline-primary btn-sm" id="exportPdf">PDF</button>
                  <button class="btn btn-outline-primary btn-sm" id="exportPptx">PowerPoint</button>
                </div>
              </div>
            </div>
          </div>
          <div class="col-12">
            <div class="card border-primary" id="frCard" style="display:none;">
              <div class="card-header bg-primary text-white d-flex align-items-center justify-content-between">
                <div>
                  <strong>⚡ FastReport Powered Reports</strong>
                  <span class="badge bg-light text-primary ms-2" id="frBadge">Detecting…</span>
                </div>
                <div>
                  <span class="small opacity-75 me-3" id="frInfo"></span>
                  <a class="btn btn-sm btn-outline-light" href="/builder">Manage Templates →</a>
                </div>
              </div>
              <div class="card-body">
                <div class="row g-3">
                  <div class="col-12 col-md-5">
                    <label class="form-label fw-bold">FastReport Template (.frx)</label>
                    <select class="form-select" id="frTemplateSelect"></select>
                    <div class="form-text" id="frTplInfo">Select a template to preview available formats.</div>
                  </div>
                  <div class="col-12 col-md-4">
                    <label class="form-label fw-bold">Export Format</label>
                    <select class="form-select" id="frFormatSelect">
                      <option value="pdf">PDF (Recommended)</option>
                      <option value="xlsx">Excel (.xlsx)</option>
                      <option value="html">HTML</option>
                      <option value="docx">Word (.docx)</option>
                      <option value="csv">CSV</option>
                      <option value="rtf">Rich Text (.rtf)</option>
                      <option value="png">PNG Image</option>
                      <option value="jpg">JPEG Image</option>
                    </select>
                    <div class="form-text">FastReport generates high-quality paginated exports.</div>
                  </div>
                  <div class="col-12 col-md-3 d-flex align-items-end">
                    <div class="d-grid gap-2 w-100">
                      <button class="btn btn-primary fw-bold" id="frGenerateBtn">
                        🚀 Generate with FastReport
                      </button>
                      <button class="btn btn-outline-secondary btn-sm" id="frRefreshBtn" type="button">
                        🔄 Refresh templates
                      </button>
                    </div>
                  </div>
                </div>
                <div class="mt-3">
                  <div class="alert alert-info py-2 mb-0 small" role="alert" id="frAlert">
                    ℹ️ When FastReport Desktop is installed, all export buttons above (CSV/Excel/PDF/PowerPoint) automatically prefer FastReport rendering with a fallback to the legacy Python renderer.
                  </div>
                </div>
              </div>
            </div>
          </div>
          <div class="col-12">
            <div id="widgets" class="row g-3"></div>
          </div>
        </div>
      </main>
      <script>
        const state = { auto: true, timer: null, templateId: null, role: 'sales', token: '' };
        function pad2(n) { return String(n).padStart(2, '0'); }
        function toIsoDate(d) { return `${d.getFullYear()}-${pad2(d.getMonth() + 1)}-${pad2(d.getDate())}`; }
        function isValidYMD(y, m, d) {
          if (!Number.isInteger(y) || !Number.isInteger(m) || !Number.isInteger(d)) return false;
          if (y < 1900 || y > 2100) return false;
          if (m < 1 || m > 12) return false;
          if (d < 1 || d > 31) return false;
          const dt = new Date(Date.UTC(y, m - 1, d));
          return dt.getUTCFullYear() === y && (dt.getUTCMonth() + 1) === m && dt.getUTCDate() === d;
        }

        function parseFlexibleDate(raw) {
          const s = String(raw || '').trim();
          if (!s) return { ok: true, iso: '' };

          const isoMatch = s.match(/^(\\d{4})-(\\d{1,2})-(\\d{1,2})$/);
          if (isoMatch) {
            const y = Number(isoMatch[1]), m = Number(isoMatch[2]), d = Number(isoMatch[3]);
            if (!isValidYMD(y, m, d)) return { ok: false, reason: 'Invalid date.' };
            return { ok: true, iso: `${y}-${pad2(m)}-${pad2(d)}` };
          }

          const parts = s.split(/[\\/\\-\\.]/).map(x => x.trim()).filter(Boolean);
          if (parts.length === 3 && parts.every(p => /^\\d+$/.test(p))) {
            const a = Number(parts[0]), b = Number(parts[1]), c = Number(parts[2]);

            if (parts[0].length === 4) {
              const y = a, m = b, d = c;
              if (!isValidYMD(y, m, d)) return { ok: false, reason: 'Invalid date.' };
              return { ok: true, iso: `${y}-${pad2(m)}-${pad2(d)}` };
            }

            let d = a, m = b, y = c;
            if (a <= 12 && b <= 12) { d = a; m = b; }
            else if (a > 12 && b <= 12) { d = a; m = b; }
            else if (b > 12 && a <= 12) { m = a; d = b; }
            else { d = a; m = b; }

            if (y < 100) y = 2000 + y;
            if (!isValidYMD(y, m, d)) return { ok: false, reason: 'Invalid date.' };
            return { ok: true, iso: `${y}-${pad2(m)}-${pad2(d)}` };
          }

          const digits = s.replace(/\D/g, '');
          if (digits.length === 8) {
            const first4 = Number(digits.slice(0, 4));
            if (first4 >= 1900 && first4 <= 2100) {
              const y = first4, m = Number(digits.slice(4, 6)), d = Number(digits.slice(6, 8));
              if (!isValidYMD(y, m, d)) return { ok: false, reason: 'Invalid date.' };
              return { ok: true, iso: `${y}-${pad2(m)}-${pad2(d)}` };
            }
            const d = Number(digits.slice(0, 2)), m = Number(digits.slice(2, 4)), y = Number(digits.slice(4, 8));
            if (!isValidYMD(y, m, d)) return { ok: false, reason: 'Invalid date.' };
            return { ok: true, iso: `${y}-${pad2(m)}-${pad2(d)}` };
          }

          return { ok: false, reason: 'Use YYYY-MM-DD, DD/MM/YYYY, or MM/DD/YYYY.' };
        }

        function setDateInputState(input, errorEl, iso, ok, message) {
          input.dataset.iso = iso || '';
          if (ok) {
            input.classList.remove('is-invalid');
            if (iso) input.classList.add('is-valid');
            else input.classList.remove('is-valid');
            if (errorEl) { errorEl.style.display = 'none'; errorEl.textContent = ''; }
          } else {
            input.classList.remove('is-valid');
            input.classList.add('is-invalid');
            if (errorEl) { errorEl.style.display = 'block'; errorEl.textContent = message || 'Invalid date.'; }
          }
        }

        function buildCalendar(popover, ctx) {
          popover.innerHTML = '';
          popover.tabIndex = -1;
          popover.setAttribute('role', 'dialog');
          popover.setAttribute('aria-label', 'Calendar');

          const head = document.createElement('div');
          head.className = 'dp-head';

          const prev = document.createElement('button');
          prev.type = 'button';
          prev.className = 'dp-btn';
          prev.textContent = '‹';
          prev.setAttribute('aria-label', 'Previous month');

          const next = document.createElement('button');
          next.type = 'button';
          next.className = 'dp-btn';
          next.textContent = '›';
          next.setAttribute('aria-label', 'Next month');

          const title = document.createElement('div');
          title.className = 'dp-title';

          head.appendChild(prev);
          head.appendChild(title);
          head.appendChild(next);
          popover.appendChild(head);

          const grid = document.createElement('div');
          grid.className = 'dp-grid';
          popover.appendChild(grid);

          function render() {
            const monthNames = ['January','February','March','April','May','June','July','August','September','October','November','December'];
            title.textContent = `${monthNames[ctx.viewMonth]} ${ctx.viewYear}`;

            grid.innerHTML = '';
            const dows = ['Mo','Tu','We','Th','Fr','Sa','Su'];
            dows.forEach(x => {
              const el = document.createElement('div');
              el.className = 'dp-dow';
              el.textContent = x;
              grid.appendChild(el);
            });

            const first = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 1));
            const firstDow = (first.getUTCDay() + 6) % 7;
            const daysInMonth = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth + 1, 0)).getUTCDate();
            const daysPrevMonth = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 0)).getUTCDate();

            const totalCells = 42;
            const selectedIso = ctx.selectedIso;
            const todayIso = toIsoDate(new Date());

            for (let i = 0; i < totalCells; i++) {
              const cellIndex = i - firstDow + 1;
              let y = ctx.viewYear, m = ctx.viewMonth + 1, d = cellIndex;
              let muted = false;
              if (cellIndex <= 0) {
                muted = true;
                const pm = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 0));
                y = pm.getUTCFullYear(); m = pm.getUTCMonth() + 1; d = daysPrevMonth + cellIndex;
              } else if (cellIndex > daysInMonth) {
                muted = true;
                const nm = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth + 1, 1));
                y = nm.getUTCFullYear(); m = nm.getUTCMonth() + 1; d = cellIndex - daysInMonth;
              }
              const iso = `${y}-${pad2(m)}-${pad2(d)}`;
              const btn = document.createElement('button');
              btn.type = 'button';
              btn.className = `dp-day${muted ? ' dp-muted' : ''}`;
              btn.textContent = String(d);
              btn.dataset.iso = iso;
              btn.setAttribute('role', 'gridcell');
              btn.setAttribute('aria-selected', iso === selectedIso ? 'true' : 'false');
              btn.setAttribute('aria-label', iso);
              if (iso === todayIso) btn.style.borderColor = 'rgba(13,110,253,.35)';

              btn.addEventListener('click', () => ctx.onSelect(iso));
              grid.appendChild(btn);
            }
          }

          prev.addEventListener('click', () => {
            const d = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 1));
            d.setUTCMonth(d.getUTCMonth() - 1);
            ctx.viewYear = d.getUTCFullYear();
            ctx.viewMonth = d.getUTCMonth();
            render();
          });
          next.addEventListener('click', () => {
            const d = new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 1));
            d.setUTCMonth(d.getUTCMonth() + 1);
            ctx.viewYear = d.getUTCFullYear();
            ctx.viewMonth = d.getUTCMonth();
            render();
          });

          popover.addEventListener('keydown', (e) => {
            if (e.key === 'Escape') { ctx.close(); return; }
            const move = (deltaDays) => {
              const base = ctx.selectedIso ? new Date(ctx.selectedIso + 'T00:00:00Z') : new Date(Date.UTC(ctx.viewYear, ctx.viewMonth, 1));
              base.setUTCDate(base.getUTCDate() + deltaDays);
              const iso = toIsoDate(new Date(base.getTime()));
              ctx.onSelect(iso, true);
            };
            if (e.key === 'ArrowLeft') { e.preventDefault(); move(-1); }
            if (e.key === 'ArrowRight') { e.preventDefault(); move(1); }
            if (e.key === 'ArrowUp') { e.preventDefault(); move(-7); }
            if (e.key === 'ArrowDown') { e.preventDefault(); move(7); }
            if (e.key === 'Enter') { e.preventDefault(); if (ctx.selectedIso) ctx.onSelect(ctx.selectedIso); }
          });

          render();
        }

        function initDateField(inputId, errorId) {
          const input = document.getElementById(inputId);
          const errorEl = document.getElementById(errorId);
          if (!input) return;

          const container = input.closest('[data-date-field]');
          if (!container) return;

          const clearBtn = container.querySelector('.date-clear');
          const openBtn = container.querySelector('.date-open');

          const popover = document.createElement('div');
          popover.className = 'date-popover';
          popover.hidden = true;
          container.appendChild(popover);

          function close() { popover.hidden = true; }
          function open(focusPopover) {
            const iso = input.dataset.iso || '';
            const base = iso ? new Date(iso + 'T00:00:00Z') : new Date();
            const ctx = {
              viewYear: base.getUTCFullYear(),
              viewMonth: base.getUTCMonth(),
              selectedIso: iso,
              onSelect: (newIso, keepOpen) => {
                input.value = newIso;
                setDateInputState(input, errorEl, newIso, true, '');
                ctx.selectedIso = newIso;
                const d = new Date(newIso + 'T00:00:00Z');
                ctx.viewYear = d.getUTCFullYear();
                ctx.viewMonth = d.getUTCMonth();
                buildCalendar(popover, ctx);
                if (!keepOpen) {
                  close();
                  input.focus();
                }
              },
              close,
            };
            buildCalendar(popover, ctx);
            popover.hidden = false;
            if (focusPopover) popover.focus();
          }

          function validateLive() {
            const raw = input.value || '';
            const digitsOnly = raw.replace(/\D/g, '');
            if (digitsOnly && digitsOnly === raw) {
              let formatted = raw;
              if ((digitsOnly.startsWith('19') || digitsOnly.startsWith('20')) && digitsOnly.length <= 8) {
                if (digitsOnly.length <= 4) formatted = digitsOnly;
                else if (digitsOnly.length <= 6) formatted = `${digitsOnly.slice(0,4)}-${digitsOnly.slice(4)}`;
                else formatted = `${digitsOnly.slice(0,4)}-${digitsOnly.slice(4,6)}-${digitsOnly.slice(6)}`;
              } else if (digitsOnly.length <= 8) {
                if (digitsOnly.length <= 2) formatted = digitsOnly;
                else if (digitsOnly.length <= 4) formatted = `${digitsOnly.slice(0,2)}-${digitsOnly.slice(2)}`;
                else formatted = `${digitsOnly.slice(0,2)}-${digitsOnly.slice(2,4)}-${digitsOnly.slice(4)}`;
              }
              if (formatted !== raw) {
                const prevLen = raw.length;
                input.value = formatted;
                try { input.setSelectionRange(formatted.length, formatted.length); } catch (e) {}
                if (formatted.length < prevLen) input.value = formatted;
              }
            }

            const parsed = parseFlexibleDate(input.value);
            if (parsed.ok) setDateInputState(input, errorEl, parsed.iso, true, '');
            else setDateInputState(input, errorEl, '', false, parsed.reason);
          }

          input.addEventListener('input', validateLive);
          input.addEventListener('blur', () => {
            const parsed = parseFlexibleDate(input.value);
            if (parsed.ok) {
              input.value = parsed.iso || '';
              setDateInputState(input, errorEl, parsed.iso, true, '');
            } else {
              setDateInputState(input, errorEl, '', false, parsed.reason);
            }
          });
          input.addEventListener('keydown', (e) => {
            if (e.key === 'Enter') {
              const parsed = parseFlexibleDate(input.value);
              if (parsed.ok) {
                input.value = parsed.iso || '';
                setDateInputState(input, errorEl, parsed.iso, true, '');
              } else {
                setDateInputState(input, errorEl, '', false, parsed.reason);
              }
              return;
            }
            if (e.key === 'ArrowDown' && e.altKey) { e.preventDefault(); open(true); return; }
            if (e.key === 'F4') { e.preventDefault(); open(true); return; }
            if (e.key === 'Escape') { close(); return; }
          });

          if (openBtn) openBtn.addEventListener('click', () => (popover.hidden ? open(false) : close()));
          if (clearBtn) clearBtn.addEventListener('click', () => {
            input.value = '';
            setDateInputState(input, errorEl, '', true, '');
            close();
          });

          document.addEventListener('click', (e) => { if (!container.contains(e.target)) close(); });
        }

        function getIsoDateValue(inputId) {
          const input = document.getElementById(inputId);
          if (!input) return '';
          const parsed = parseFlexibleDate(input.value);
          if (parsed.ok) return parsed.iso || '';
          return '';
        }

        function getIsoDateValueOrWarn(inputId, errorId) {
          const input = document.getElementById(inputId);
          const errorEl = document.getElementById(errorId);
          if (!input) return '';
          const parsed = parseFlexibleDate(input.value);
          if (parsed.ok) return parsed.iso || '';
          if ((input.value || '').trim()) {
            setDateInputState(input, errorEl, '', false, parsed.reason);
            input.focus();
            return null;
          }
          return '';
        }

        function setDateFieldValue(inputId, errorId, value) {
          const input = document.getElementById(inputId);
          const errorEl = document.getElementById(errorId);
          if (!input) return;
          const parsed = parseFlexibleDate(value);
          if (parsed.ok) {
            input.value = parsed.iso || '';
            setDateInputState(input, errorEl, parsed.iso, true, '');
          } else {
            input.value = String(value || '');
            setDateInputState(input, errorEl, '', false, parsed.reason);
          }
        }

        function qs() {
          const p = new URLSearchParams(window.location.search);
          return {
            from_date: p.get('from_date') || '',
            to_date: p.get('to_date') || '',
            status: p.get('status') || 'ALL',
            template_id: p.get('template_id') || ''
          };
        }
        function headers() {
          const h = { 'Content-Type': 'application/json' };
          if (state.role) h['X-User-Role'] = state.role;
          if (state.token) h['X-API-Key'] = state.token;
          return h;
        }
        async function loadTemplates() {
          const res = await fetch('/api/templates', { headers: headers() });
          const data = await res.json();
          const sel = document.getElementById('templateSelect');
          sel.innerHTML = '';
          data.items.forEach(t => {
            const opt = document.createElement('option');
            opt.value = t.id;
            opt.textContent = t.name;
            sel.appendChild(opt);
          });
          const q = qs();
          if (q.template_id && [...sel.options].some(o => o.value == q.template_id)) {
            sel.value = q.template_id;
          } else {
            const defaultOpt = [...sel.options].find(o => o.textContent === 'Default Dashboard') || sel.options[0];
            if (defaultOpt) sel.value = defaultOpt.value;
          }
          state.templateId = sel.value;
        }
        function renderKpi(title, value) {
          const col = document.createElement('div');
          col.className = 'col-12 col-md-3';
          col.innerHTML = `<div class="card"><div class="card-body"><div class="text-muted small">${title}</div><div class="fs-4 fw-bold">${value}</div></div></div>`;
          return col;
        }
        function renderTable(title, rows) {
          const col = document.createElement('div');
          col.className = 'col-12';
          const head = `<thead><tr><th>Invoice</th><th>Date</th><th>POS</th><th>Mode</th><th>Total</th><th>Status</th><th style="width: 110px;">Action</th></tr></thead>`;
          const body = rows.map(r => {
            const inv = r.invoice_number;
            const statusTd = `<td data-status-for="${inv}">${r.sync_status}</td>`;
            const canRetry = ['PENDING', 'FAILED'].includes((r.sync_status || '').toUpperCase());
            const btn = canRetry
              ? `<button type="button" class="btn btn-sm btn-outline-primary retry-btn" data-invoice="${inv}">Retry</button>`
              : `<span class="text-muted small">—</span>`;
            return `<tr data-invoice-row="${inv}"><td>${inv}</td><td>${r.datetime}</td><td>${r.pos_id}</td><td>${r.payment_mode}</td><td>${r.total_amount.toFixed(2)}</td>${statusTd}<td>${btn}</td></tr>`;
          }).join('');
          col.innerHTML = `<div class="card"><div class="card-header fw-bold">${title}</div><div class="card-body p-0"><div class="table-responsive"><table class="table table-sm mb-0">${head}<tbody>${body}</tbody></table></div></div></div>`;
          return col;
        }
        function renderChart(title, dailySales) {
          const col = document.createElement('div');
          col.className = 'col-12';
          const id = 'chart_' + Math.random().toString(36).slice(2);
          col.innerHTML = `<div class="card"><div class="card-header fw-bold">${title}</div><div class="card-body"><div id="${id}" style="height: 320px;"></div></div></div>`;
          setTimeout(() => {
            if (Array.isArray(dailySales) && dailySales.length && dailySales[0].date !== undefined) {
              const x = dailySales.map(d => d.date);
              const y = dailySales.map(d => d.amount);
              Plotly.newPlot(id, [{ x, y, type: 'scatter', mode: 'lines+markers' }], { margin: { t: 10, r: 10, l: 40, b: 40 } }, { displayModeBar: false, responsive: true });
              return;
            }
            if (Array.isArray(dailySales) && dailySales.length && dailySales[0].status !== undefined) {
              const labels = dailySales.map(d => d.status || 'UNKNOWN');
              const values = dailySales.map(d => d.count || 0);
              Plotly.newPlot(id, [{ labels, values, type: 'pie' }], { margin: { t: 10, r: 10, l: 10, b: 10 } }, { displayModeBar: false, responsive: true });
              return;
            }
            Plotly.newPlot(id, [], { margin: { t: 10, r: 10, l: 10, b: 10 } }, { displayModeBar: false, responsive: true });
          }, 0);
          return col;
        }
        async function loadDashboard() {
          const sel = document.getElementById('templateSelect');
          state.templateId = sel.value;
          const from_date = getIsoDateValueOrWarn('fromDate', 'fromDateError');
          if (from_date === null) return;
          const to_date = getIsoDateValueOrWarn('toDate', 'toDateError');
          if (to_date === null) return;
          const status = document.getElementById('status').value;
          const url = `/api/dashboard?template_id=${encodeURIComponent(state.templateId)}&from_date=${encodeURIComponent(from_date)}&to_date=${encodeURIComponent(to_date)}&status=${encodeURIComponent(status)}`;
          const res = await fetch(url, { headers: headers() });
          const data = await res.json();
          const w = document.getElementById('widgets');
          w.innerHTML = '';
          if (!data.widgets || data.widgets.length === 0) {
            w.innerHTML = `
              <div class="col-12">
                <div class="card border-warning">
                  <div class="card-header fw-bold bg-warning-subtle">Custom template (no preview widgets defined).</div>
                  <div class="card-body">
                    <p class="mb-2">Edit the layout visually in <a class="fw-bold" href="/builder">FastReport Studio</a> using the Designer. Exports always use FastReport templates regardless of dashboard preview widgets.</p>
                    <p class="mb-0 text-muted small">Or switch to the built-in <b>Default Dashboard</b> template from the dropdown above for a dashboard preview.</p>
                  </div>
                </div>
              </div>`;
            return;
          }
          data.widgets.forEach(widget => {
            if (widget.type === 'kpi') w.appendChild(renderKpi(widget.title, widget.value));
            if (widget.type === 'chart') w.appendChild(renderChart(widget.title, widget.value));
            if (widget.type === 'table') w.appendChild(renderTable(widget.title, widget.value));
          });
        }
        function toast(message, kind) {
          const id = 'toast_' + Math.random().toString(36).slice(2);
          const bg = kind === 'success' ? 'text-bg-success' : (kind === 'error' ? 'text-bg-danger' : 'text-bg-secondary');
          const el = document.createElement('div');
          el.innerHTML = `
            <div class="toast-container position-fixed bottom-0 end-0 p-3">
              <div id="${id}" class="toast align-items-center ${bg} border-0" role="alert" aria-live="assertive" aria-atomic="true">
                <div class="d-flex">
                  <div class="toast-body">${message}</div>
                  <button type="button" class="btn-close btn-close-white me-2 m-auto" data-bs-dismiss="toast" aria-label="Close"></button>
                </div>
              </div>
            </div>`;
          document.body.appendChild(el);
          const toastEl = document.getElementById(id);
          const t = new bootstrap.Toast(toastEl, { delay: 4500 });
          t.show();
          toastEl.addEventListener('hidden.bs.toast', () => el.remove());
        }
        async function retryInvoice(invoiceNumber, buttonEl) {
          if (!invoiceNumber || !buttonEl) return;
          const now = Date.now();
          const last = parseInt(buttonEl.getAttribute('data-last-click') || '0', 10);
          if (now - last < 15000) {
            toast('Please wait before retrying again.', 'info');
            return;
          }
          buttonEl.setAttribute('data-last-click', String(now));
          const originalHtml = buttonEl.innerHTML;
          buttonEl.disabled = true;
          buttonEl.innerHTML = `<span class="spinner-border spinner-border-sm me-1" role="status" aria-hidden="true"></span>Retrying`;
          const statusCell = document.querySelector(`td[data-status-for="${invoiceNumber}"]`);
          if (statusCell) statusCell.textContent = 'RETRYING';
          try {
            const res = await fetch(`/api/invoices/${encodeURIComponent(invoiceNumber)}/retry`, { method: 'POST', headers: headers() });
            const data = await res.json();
            if (!res.ok) {
              const msg = data && (data.detail || data.message) ? (data.detail || data.message) : 'Retry failed.';
              toast(msg, 'error');
              if (statusCell) statusCell.textContent = (data && data.sync_status) ? data.sync_status : 'FAILED';
              buttonEl.disabled = false;
              buttonEl.innerHTML = originalHtml;
              return;
            }
            const newStatus = (data && data.sync_status) ? data.sync_status : 'PENDING';
            if (statusCell) statusCell.textContent = newStatus;
            if (newStatus.toUpperCase() === 'SYNCED') {
              toast(`Invoice ${invoiceNumber} uploaded successfully.`, 'success');
              buttonEl.outerHTML = `<span class="text-muted small">—</span>`;
            } else if (newStatus.toUpperCase() === 'FAILED') {
              toast(data && data.message ? data.message : `Invoice ${invoiceNumber} failed.`, 'error');
              buttonEl.disabled = false;
              buttonEl.innerHTML = originalHtml;
            } else {
              toast(data && data.message ? data.message : `Invoice ${invoiceNumber} is still pending.`, 'info');
              buttonEl.disabled = false;
              buttonEl.innerHTML = originalHtml;
            }
          } catch (e) {
            toast(`Network error: ${e}`, 'error');
            if (statusCell) statusCell.textContent = 'FAILED';
            buttonEl.disabled = false;
            buttonEl.innerHTML = originalHtml;
          }
        }
        function setAuto(on) {
          state.auto = on;
          const btn = document.getElementById('autoBtn');
          btn.textContent = on ? 'Auto' : 'Manual';
          if (state.timer) clearInterval(state.timer);
          if (on) state.timer = setInterval(loadDashboard, 15000);
        }
        function exportFmt(fmt) {
          const from_date = getIsoDateValueOrWarn('fromDate', 'fromDateError');
          if (from_date === null) return;
          const to_date = getIsoDateValueOrWarn('toDate', 'toDateError');
          if (to_date === null) return;
          const status = document.getElementById('status').value;
          const url = `/export/${fmt}?template_id=${encodeURIComponent(state.templateId)}&from_date=${encodeURIComponent(from_date)}&to_date=${encodeURIComponent(to_date)}&status=${encodeURIComponent(status)}`;
          window.location.href = url;
        }
        async function init() {
          const q = qs();
          initDateField('fromDate', 'fromDateError');
          initDateField('toDate', 'toDateError');
          setDateFieldValue('fromDate', 'fromDateError', q.from_date);
          setDateFieldValue('toDate', 'toDateError', q.to_date);
          document.getElementById('status').value = q.status;
          await loadTemplates();
          document.getElementById('applyBtn').addEventListener('click', loadDashboard);
          document.getElementById('autoBtn').addEventListener('click', () => setAuto(!state.auto));
          document.getElementById('templateSelect').addEventListener('change', loadDashboard);
          document.getElementById('exportCsv').addEventListener('click', () => exportFmt('csv'));
          document.getElementById('exportXlsx').addEventListener('click', () => exportFmt('xlsx'));
          document.getElementById('exportPdf').addEventListener('click', () => exportFmt('pdf'));
          document.getElementById('exportPptx').addEventListener('click', () => exportFmt('pptx'));
          document.addEventListener('click', (ev) => {
            const t = ev.target;
            if (!t || !t.classList) return;
            if (t.classList.contains('retry-btn')) {
              ev.preventDefault();
              retryInvoice(t.getAttribute('data-invoice'), t);
            }
          });

          // ---- FastReport integration ----
          function headers() {
            const h = {};
            if (state.token) h['X-API-Key'] = state.token;
            if (state.role)  h['X-User-Role'] = state.role;
            return h;
          }
          function currentFilters() {
            const from_ = document.getElementById('fromDate').dataset.iso || '';
            const to = document.getElementById('toDate').dataset.iso || '';
            const status = document.getElementById('status').value;
            return { from_date: from_, to_date: to, status: status };
          }
          function downloadURL(template_name, fmt) {
            const f = currentFilters();
            const params = new URLSearchParams();
            params.set('fmt', fmt);
            if (f.from_date) params.set('from_date', f.from_date);
            if (f.to_date)   params.set('to_date',   f.to_date);
            if (f.status)    params.set('status',    f.status);
            return `/api/fastreports/export/${encodeURIComponent(template_name)}?${params.toString()}`;
          }
          async function loadFastReportStatus() {
            const card = document.getElementById('frCard');
            const badge = document.getElementById('frBadge');
            const info = document.getElementById('frInfo');
            const tplSel = document.getElementById('frTemplateSelect');
            const tplInfo = document.getElementById('frTplInfo');
            if (!card) return;
            try {
              const res = await fetch('/api/fastreports/status', { headers: headers() });
              if (!res.ok) { card.style.display = 'none'; return; }
              const data = await res.json();
              card.style.display = '';
              if (data.available) {
                badge.className = 'badge bg-success text-white ms-2';
                badge.textContent = 'ACTIVE';
                info.textContent = `Builder: ${(data.builder_exe || 'n/a').split('\\').pop()}`;
              } else {
                badge.className = 'badge bg-warning text-dark ms-2';
                badge.textContent = 'NOT INSTALLED';
                info.textContent = 'Set FASTREPORT_DESKTOP_DIR or install FastReport Desktop.';
                document.getElementById('frAlert').className = 'alert alert-warning py-2 mb-0 small';
                document.getElementById('frAlert').textContent = '⚠️ FastReport Desktop not detected. Legacy Python-native exports still work (reportlab/openpyxl). Install FastReport for higher-quality PDF/XLSX/DOCX output.';
              }
              tplSel.innerHTML = '';
              const templates = data.templates || [];
              if (!templates.length) {
                tplSel.innerHTML = '<option value="">(no .frx templates found)</option>';
                tplInfo.textContent = 'No .frx templates in exports/templates_frx/';
              } else {
                templates.forEach((t) => {
                  const opt = document.createElement('option');
                  opt.value = t.name;
                  const sizeKB = Math.round(t.size_bytes / 1024);
                  opt.textContent = t.name + ' (' + sizeKB + ' KB, ' + t.modified_at + ')';
                  tplSel.appendChild(opt);
                });
                tplInfo.textContent = templates.length + ' .frx template(s) loaded from exports/templates_frx/';
              }
            } catch (e) {
              card.style.display = 'none';
            }
          }
          document.getElementById('frGenerateBtn')?.addEventListener('click', () => {
            const tpl = document.getElementById('frTemplateSelect').value;
            const fmt = document.getElementById('frFormatSelect').value;
            if (!tpl) { alert('Please select a FastReport template.'); return; }
            const btn = document.getElementById('frGenerateBtn');
            const oldText = btn.textContent;
            btn.disabled = true;
            btn.textContent = '⏳ Generating…';
            try {
              const form = document.createElement('form');
              form.method = 'POST';
              form.action = downloadURL(tpl, fmt);
              const token = state.token;
              if (token) {
                const tok = document.createElement('input');
                tok.type = 'hidden'; tok.name = 'x_api_key'; tok.value = token;
                form.appendChild(tok);
              }
              document.body.appendChild(form);
              form.submit();
              document.body.removeChild(form);
            } finally {
              setTimeout(() => { btn.disabled = false; btn.textContent = oldText; }, 1500);
            }
          });
          document.getElementById('frRefreshBtn')?.addEventListener('click', loadFastReportStatus);

          setAuto(true);
          await loadDashboard();
          await loadFastReportStatus();
        }
        init();
      </script>
    </body>
    </html>
    """


def _render_builder_html() -> str:
    return """
    <!doctype html>
    <html lang="en">
    <head>
      <meta charset="utf-8"/>
      <meta name="viewport" content="width=device-width, initial-scale=1"/>
      <title>FastReport Template Studio</title>
      <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
    </head>
    <body class="bg-light">
      <nav class="navbar navbar-expand-lg navbar-dark bg-primary">
        <div class="container-fluid">
          <a class="navbar-brand" href="/dashboard">Reporting Portal · FastReport Studio</a>
          <div class="d-flex gap-2">
            <a class="btn btn-outline-light btn-sm" href="/dashboard">Dashboard</a>
            <a class="btn btn-outline-light btn-sm" href="/schedules">Schedules</a>
            <a class="btn btn-outline-light btn-sm" href="/lookup">Lookup</a>
          </div>
        </div>
      </nav>
      <main class="container py-3">
        <div class="row g-3 mb-3">
          <div class="col-12 col-lg-6">
            <div class="card border-primary h-100">
              <div class="card-header bg-primary text-white fw-bold">
                ⚡ FastReport Desktop Status
              </div>
              <div class="card-body">
                <div class="d-flex align-items-center gap-2 mb-2">
                  <span class="spinner-border spinner-border-sm text-primary" id="frSpinner" role="status"></span>
                  <span class="fw-bold" id="frStatusText">Detecting FastReport installation…</span>
                </div>
                <div id="frDetails" class="small text-muted"></div>
                <div class="mt-3 d-flex flex-wrap gap-2">
                  <button class="btn btn-primary" id="openDesignerBlankBtn">
                    🎨 Open Designer (Blank)
                  </button>
                  <button class="btn btn-outline-secondary btn-sm" id="refreshStatusBtn" type="button">
                    🔄 Refresh
                  </button>
                </div>
                <div class="mt-3 small">
                  <div class="fw-bold mb-1">Tips:</div>
                  <ul class="mb-0 ps-3">
                    <li>FastReport Designer edits <code>.frx</code> XML templates.</li>
                    <li>Templates are stored in <code>exports/templates_frx/</code>.</li>
                    <li>Data fields: <code>[Data.report_title]</code>, <code>[Data.invoices]</code>, <code>[Data.total_amount]</code>, etc.</li>
                    <li>Design changes are saved locally to the <code>.frx</code> file instantly.</li>
                  </ul>
                </div>
              </div>
            </div>
          </div>
          <div class="col-12 col-lg-6">
            <div class="card h-100">
              <div class="card-header fw-bold">Quick Export Test</div>
              <div class="card-body">
                <div class="row g-2 mb-3">
                  <div class="col-7">
                    <label class="form-label fw-bold">Template</label>
                    <select class="form-select" id="quickTemplateSelect"></select>
                  </div>
                  <div class="col-5">
                    <label class="form-label fw-bold">Format</label>
                    <select class="form-select" id="quickFormatSelect">
                      <option value="pdf">PDF</option>
                      <option value="xlsx">Excel (.xlsx)</option>
                      <option value="html">HTML</option>
                      <option value="docx">Word (.docx)</option>
                      <option value="csv">CSV</option>
                      <option value="rtf">RTF</option>
                      <option value="png">PNG</option>
                    </select>
                  </div>
                </div>
                <button class="btn btn-success w-100 fw-bold" id="quickExportBtn">
                  🚀 Render &amp; Download Now
                </button>
                <div class="form-text mt-2">Renders the selected template with live invoice data from the last 30 days.</div>
              </div>
            </div>
          </div>
        </div>
        <div class="row g-3">
          <div class="col-12 col-lg-5">
            <div class="card">
              <div class="card-header fw-bold d-flex align-items-center justify-content-between">
                <span>FastReport Templates (.frx)</span>
                <button class="btn btn-outline-primary btn-sm" id="refreshTemplatesBtn" type="button">🔄</button>
              </div>
              <div class="card-body">
                <div class="small text-muted mb-2" id="tplCountInfo">Loading templates…</div>
                <div class="table-responsive" style="max-height: 380px; overflow-y: auto;">
                  <table class="table table-sm table-hover mb-0" id="frxTable">
                    <thead class="table-light sticky-top">
                      <tr>
                        <th>Template</th>
                        <th class="text-end">Size</th>
                        <th>Modified</th>
                        <th>Actions</th>
                      </tr>
                    </thead>
                    <tbody id="frxBody"><tr><td colspan="4" class="text-center text-muted py-3">(loading…)</td></tr></tbody>
                  </table>
                </div>
              </div>
            </div>
            <div class="card mt-3">
              <div class="card-header fw-bold">Schedule Templates (DB)</div>
              <div class="card-body small text-muted">
                <div class="mb-2">The database stores which template is selected per Schedule entry. Edit Schedule templates on the <a href="/schedules" class="fw-bold">Schedules</a> page.</div>
                <select class="form-select form-select-sm" id="dbTemplateSelect" size="6" aria-label="DB Templates list"></select>
              </div>
            </div>
          </div>
          <div class="col-12 col-lg-7">
            <div class="card border-info h-100">
              <div class="card-header bg-info-subtle fw-bold d-flex align-items-center justify-content-between">
                <span>🛠️ Template: <span id="selectedTplName" class="text-info">— select one on the left —</span></span>
                <span class="badge bg-info text-white" id="selectedTplBadge">idle</span>
              </div>
              <div class="card-body">
                <div class="alert alert-info small py-2 mb-3">
                  <strong>Workflow:</strong> Click <em>Open in Designer</em> to edit the template layout using FastReport Designer. After saving the <code>.frx</code> file, click <em>Preview / Export</em> to render it against live invoice data.
                </div>
                <div class="row g-3 mb-3">
                  <div class="col-6 col-md-4">
                    <label class="form-label fw-bold">File (.frx)</label>
                    <input type="text" class="form-control form-control-sm" id="tplFile" readonly placeholder="—"/>
                  </div>
                  <div class="col-6 col-md-4">
                    <label class="form-label fw-bold">Size</label>
                    <input type="text" class="form-control form-control-sm" id="tplSize" readonly placeholder="—"/>
                  </div>
                  <div class="col-12 col-md-4">
                    <label class="form-label fw-bold">Modified</label>
                    <input type="text" class="form-control form-control-sm" id="tplModified" readonly placeholder="—"/>
                  </div>
                </div>
                <div class="d-flex flex-wrap gap-2 mb-3">
                  <button class="btn btn-primary" id="openInDesignerBtn" disabled>🎨 Open in Designer</button>
                  <button class="btn btn-success" id="renderPdfBtn" disabled>📄 Preview PDF</button>
                  <button class="btn btn-success" id="renderXlsxBtn" disabled>📊 Preview XLSX</button>
                  <button class="btn btn-success" id="renderHtmlBtn" disabled>🌐 Preview HTML</button>
                  <button class="btn btn-outline-success" id="renderAllBtn" disabled>📦 Render All</button>
                </div>
                <div>
                  <label class="form-label fw-bold">Quick Custom Export</label>
                  <div class="input-group">
                    <select class="form-select" id="customFmtSelect">
                      <option value="pdf">PDF</option>
                      <option value="xlsx">Excel XLSX</option>
                      <option value="html">HTML</option>
                      <option value="docx">Word DOCX</option>
                      <option value="csv">CSV</option>
                      <option value="rtf">Rich Text (RTF)</option>
                      <option value="png">Image PNG</option>
                      <option value="jpg">Image JPG</option>
                    </select>
                    <button class="btn btn-outline-primary" id="customExportBtn" disabled>Download</button>
                  </div>
                </div>
              </div>
            </div>
          </div>
        </div>
      </main>
      <script>
        const state = { token: '', role: 'admin', selectedTpl: null };
        function headers() {
          const h = { 'Content-Type': 'application/json', 'X-User-Role': state.role };
          if (state.token) h['X-API-Key'] = state.token;
          return h;
        }
        function postHeaders() {
          const h = { 'X-User-Role': state.role };
          if (state.token) h['X-API-Key'] = state.token;
          return h;
        }
        function downloadViaForm(url, bodyFields) {
          const form = document.createElement('form');
          form.method = 'POST';
          form.action = url;
          Object.entries(bodyFields || {}).forEach(([k, v]) => {
            const inp = document.createElement('input');
            inp.type = 'hidden'; inp.name = k; inp.value = v;
            form.appendChild(inp);
          });
          document.body.appendChild(form);
          form.submit();
          document.body.removeChild(form);
        }
        async function loadDBTemplates() {
          try {
            const res = await fetch('/api/templates', { headers: headers() });
            if (!res.ok) return;
            const data = await res.json();
            const sel = document.getElementById('dbTemplateSelect');
            sel.innerHTML = '';
            (data.items || []).forEach(t => {
              const opt = document.createElement('option');
              opt.value = t.id;
              opt.textContent = t.name + (t.description ? ' — ' + t.description : '');
              sel.appendChild(opt);
            });
          } catch (_) {}
        }
        function renderTplBody(templates) {
          const body = document.getElementById('frxBody');
          const countInfo = document.getElementById('tplCountInfo');
          if (!templates || !templates.length) {
            body.innerHTML = '<tr><td colspan="4" class="text-center text-muted py-3">No .frx templates found in exports/templates_frx/</td></tr>';
            if (countInfo) countInfo.textContent = '0 templates available.';
            return;
          }
          countInfo.textContent = templates.length + ' .frx template(s) — click a row to select.';
          body.innerHTML = '';
          templates.forEach((t, idx) => {
            const tr = document.createElement('tr');
            tr.style.cursor = 'pointer';
            tr.dataset.name = t.name;
            const sizeKB = Math.round(t.size_bytes / 1024);
            const sizeStr = sizeKB < 1 ? (t.size_bytes + ' B') : (sizeKB + ' KB');
            tr.innerHTML = `<td class="fw-semibold">${t.name}</td><td class="text-end">${sizeStr}</td><td>${t.modified_at || '—'}</td><td><button class="btn btn-sm btn-outline-primary act-open" data-name="${t.name}">🎨</button> <button class="btn btn-sm btn-outline-success act-pdf" data-name="${t.name}">PDF</button></td>`;
            tr.addEventListener('click', (ev) => {
              if (ev.target && ev.target.classList && (ev.target.classList.contains('act-open') || ev.target.classList.contains('act-pdf'))) return;
              selectTpl(t);
            });
            tr.querySelector('.act-open').addEventListener('click', (ev) => { ev.stopPropagation(); openInDesigner(t.name); });
            tr.querySelector('.act-pdf').addEventListener('click', (ev) => { ev.stopPropagation(); renderTpl(t.name, 'pdf'); });
            body.appendChild(tr);
          });
        }
        function selectTpl(t) {
          state.selectedTpl = t;
          const rows = document.querySelectorAll('#frxBody tr');
          rows.forEach(r => r.classList.remove('table-primary'));
          const row = document.querySelector('#frxBody tr[data-name="' + t.name + '"]');
          if (row) row.classList.add('table-primary');
          document.getElementById('selectedTplName').textContent = t.name;
          document.getElementById('selectedTplBadge').textContent = 'selected';
          document.getElementById('selectedTplBadge').className = 'badge bg-success text-white';
          document.getElementById('tplFile').value = t.name + '.frx';
          document.getElementById('tplSize').value = Math.round(t.size_bytes / 1024) + ' KB (' + t.size_bytes + ' bytes)';
          document.getElementById('tplModified').value = t.modified_at || '';
          ['openInDesignerBtn','renderPdfBtn','renderXlsxBtn','renderHtmlBtn','renderAllBtn','customExportBtn'].forEach(id => {
            const el = document.getElementById(id);
            if (el) el.disabled = false;
          });
        }
        async function refreshStatusAndTemplates() {
          const spinner = document.getElementById('frSpinner');
          const statusText = document.getElementById('frStatusText');
          const details = document.getElementById('frDetails');
          if (spinner) spinner.classList.remove('d-none');
          try {
            const res = await fetch('/api/fastreports/status', { headers: headers() });
            if (!res.ok) throw new Error('status fetch failed');
            const data = await res.json();
            if (spinner) spinner.classList.add('d-none');
            if (data.available) {
              statusText.textContent = '✅ FastReport Desktop ACTIVE';
              statusText.className = 'fw-bold text-success';
              const parts = [];
              if (data.builder_exe)   parts.push('Builder: <code>' + data.builder_exe.split('\\').pop() + '</code>');
              if (data.designer_exe)  parts.push('Designer: <code>' + data.designer_exe.split('\\').pop() + '</code>');
              if (data.templates_dir) parts.push('Templates: <code>' + data.templates_dir + '</code>');
              details.innerHTML = parts.join('<br>') || '';
            } else {
              statusText.textContent = '⚠️ FastReport Desktop NOT INSTALLED';
              statusText.className = 'fw-bold text-warning';
              details.innerHTML = 'Legacy exports (reportlab/openpyxl) still work. Install FastReport Desktop and/or set <code>FASTREPORT_DESKTOP_DIR</code> env var to enable FRX rendering.';
            }
            renderTplBody(data.templates || []);
            const quickSel = document.getElementById('quickTemplateSelect');
            if (quickSel) {
              quickSel.innerHTML = '';
              (data.templates || []).forEach(t => {
                const opt = document.createElement('option');
                opt.value = t.name;
                opt.textContent = t.name;
                quickSel.appendChild(opt);
              });
            }
          } catch (e) {
            if (spinner) spinner.classList.add('d-none');
            statusText.textContent = '❌ Could not reach /api/fastreports/status';
            statusText.className = 'fw-bold text-danger';
          }
        }
        function openInDesigner(templateName) {
          if (!confirm('Open FastReport Designer for: ' + templateName + ' ?')) return;
          fetch('/api/fastreports/designer/open', {
            method: 'POST',
            headers: headers(),
            body: JSON.stringify({ template_name: templateName })
          }).then(async (r) => {
            const data = await r.json().catch(() => ({}));
            if (!r.ok) { alert('Failed to open designer: ' + (data.detail || r.statusText)); return; }
            const badge = document.getElementById('selectedTplBadge');
            if (badge) { badge.textContent = 'designer launched'; badge.className = 'badge bg-primary text-white'; }
            setTimeout(() => { if (badge && state.selectedTpl) { badge.textContent = 'selected'; badge.className = 'badge bg-success text-white'; } }, 2500);
          }).catch(err => alert('Designer error: ' + err.message));
        }
        function renderTpl(templateName, fmt) {
          const url = '/api/fastreports/export/' + encodeURIComponent(templateName) + '?fmt=' + encodeURIComponent(fmt || 'pdf');
          const body = {};
          if (state.token) body.x_api_key = state.token;
          downloadViaForm(url, body);
        }
        document.getElementById('openDesignerBlankBtn')?.addEventListener('click', () => openInDesigner(''));
        document.getElementById('refreshStatusBtn')?.addEventListener('click', refreshStatusAndTemplates);
        document.getElementById('refreshTemplatesBtn')?.addEventListener('click', refreshStatusAndTemplates);
        document.getElementById('openInDesignerBtn')?.addEventListener('click', () => {
          if (state.selectedTpl) openInDesigner(state.selectedTpl.name);
        });
        document.getElementById('renderPdfBtn')?.addEventListener('click', () => { if (state.selectedTpl) renderTpl(state.selectedTpl.name, 'pdf'); });
        document.getElementById('renderXlsxBtn')?.addEventListener('click', () => { if (state.selectedTpl) renderTpl(state.selectedTpl.name, 'xlsx'); });
        document.getElementById('renderHtmlBtn')?.addEventListener('click', () => { if (state.selectedTpl) renderTpl(state.selectedTpl.name, 'html'); });
        document.getElementById('renderAllBtn')?.addEventListener('click', () => {
          if (!state.selectedTpl) return;
          ['pdf','xlsx','docx','html'].forEach((fmt, i) => {
            setTimeout(() => renderTpl(state.selectedTpl.name, fmt), i * 400);
          });
        });
        document.getElementById('customExportBtn')?.addEventListener('click', () => {
          if (!state.selectedTpl) return;
          const fmt = document.getElementById('customFmtSelect').value;
          renderTpl(state.selectedTpl.name, fmt);
        });
        document.getElementById('quickExportBtn')?.addEventListener('click', () => {
          const tpl = document.getElementById('quickTemplateSelect').value;
          const fmt = document.getElementById('quickFormatSelect').value;
          if (!tpl) { alert('Select a template first.'); return; }
          renderTpl(tpl, fmt);
        });
        refreshStatusAndTemplates();
        loadDBTemplates();
      </script>
    </body>
    </html>
    """


def _render_schedules_html() -> str:
    return """
    <!doctype html>
    <html lang="en">
    <head>
      <meta charset="utf-8"/>
      <meta name="viewport" content="width=device-width, initial-scale=1"/>
      <title>Report Schedules</title>
      <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
    </head>
    <body class="bg-light">
      <nav class="navbar navbar-expand-lg navbar-dark bg-primary">
        <div class="container-fluid">
          <a class="navbar-brand" href="/dashboard">Reporting Portal</a>
          <div class="d-flex gap-2">
            <a class="btn btn-outline-light btn-sm" href="/dashboard">Dashboard</a>
            <a class="btn btn-outline-light btn-sm" href="/builder">FastReport Studio</a>
            <a class="btn btn-outline-light btn-sm" href="/lookup">Lookup</a>
          </div>
        </div>
      </nav>
      <main class="container py-3">
        <div class="card border-info mb-3" id="frStatusCard">
          <div class="card-header bg-info-subtle fw-bold">⚡ FastReport Scheduler Note</div>
          <div class="card-body py-2 small">
            <div id="frStatusLine">Detecting FastReport…</div>
            <div class="mt-1 text-muted">Schedules always attempt FastReport first. DOCX/RTF/PNG/JPG require FastReport Desktop; PDF/XLSX/CSV will fallback to legacy renderers if FastReport is unavailable.</div>
          </div>
        </div>
        <div class="card">
          <div class="card-header fw-bold">Scheduled Reports</div>
          <div class="card-body">
            <div class="row g-2 align-items-end">
              <div class="col-12 col-md-4">
                <label class="form-label">Template</label>
                <select class="form-select" id="templateSelect"></select>
              </div>
              <div class="col-6 col-md-2">
                <label class="form-label">Interval (min)</label>
                <input class="form-control" type="number" id="interval" value="60"/>
              </div>
              <div class="col-6 col-md-2">
                <label class="form-label">Format</label>
                <select class="form-select" id="format">
                  <optgroup label="FastReport + Legacy Fallback">
                    <option value="pdf" selected>PDF</option>
                    <option value="xlsx">Excel (.xlsx)</option>
                    <option value="csv">CSV</option>
                  </optgroup>
                  <optgroup label="FastReport Desktop Only">
                    <option value="docx">Word (.docx)</option>
                    <option value="html">HTML</option>
                    <option value="rtf">RTF</option>
                    <option value="png">PNG Image</option>
                    <option value="jpg">JPG Image</option>
                  </optgroup>
                  <optgroup label="Legacy Only">
                    <option value="pptx">PowerPoint (.pptx)</option>
                  </optgroup>
                </select>
              </div>
              <div class="col-12 col-md-4">
                <label class="form-label">Recipients (comma)</label>
                <input class="form-control" id="recipients" placeholder="a@b.com,c@d.com"/>
              </div>
              <div class="col-12">
                <button class="btn btn-primary" id="createBtn">Create Schedule</button>
              </div>
            </div>
            <hr/>
            <div class="table-responsive">
              <table class="table table-sm" id="tbl">
                <thead><tr><th>ID</th><th>Template</th><th>Interval</th><th>Format</th><th>Enabled</th><th>Last Run</th><th></th></tr></thead>
                <tbody></tbody>
              </table>
            </div>
          </div>
        </div>
      </main>
      <script>
        const state = { role: 'admin', token: '' };
        function headers() {
          const h = { 'Content-Type': 'application/json', 'X-User-Role': state.role };
          if (state.token) h['X-API-Key'] = state.token;
          return h;
        }
        async function loadTemplates() {
          const res = await fetch('/api/templates', { headers: headers() });
          const data = await res.json();
          const sel = document.getElementById('templateSelect');
          sel.innerHTML = '';
          data.items.forEach(t => {
            const opt = document.createElement('option');
            opt.value = t.id;
            opt.textContent = t.name;
            sel.appendChild(opt);
          });
        }
        async function loadSchedules() {
          const res = await fetch('/api/schedules', { headers: headers() });
          const data = await res.json();
          const body = document.querySelector('#tbl tbody');
          body.innerHTML = '';
          data.items.forEach(s => {
            const tr = document.createElement('tr');
            tr.innerHTML = `<td>${s.id}</td><td>${s.template_name}</td><td>${s.interval_minutes}</td><td>${s.export_format}</td><td>${s.enabled}</td><td>${s.last_run_at || ''}</td>
              <td><button class="btn btn-sm btn-outline-secondary">Toggle</button></td>`;
            tr.querySelector('button').addEventListener('click', async () => {
              await fetch(`/api/schedules/${s.id}`, { method:'PUT', headers: headers(), body: JSON.stringify({ enabled: !s.enabled }) });
              await loadSchedules();
            });
            body.appendChild(tr);
          });
        }
        async function createSchedule() {
          const payload = {
            template_id: parseInt(document.getElementById('templateSelect').value, 10),
            interval_minutes: parseInt(document.getElementById('interval').value, 10),
            export_format: document.getElementById('format').value,
            recipients: document.getElementById('recipients').value.split(',').map(x => x.trim()).filter(Boolean),
            enabled: true
          };
          await fetch('/api/schedules', { method:'POST', headers: headers(), body: JSON.stringify(payload) });
          await loadSchedules();
        }
        document.getElementById('createBtn').addEventListener('click', createSchedule);
        async function loadFRStatus() {
          const el = document.getElementById('frStatusLine');
          if (!el) return;
          try {
            const res = await fetch('/api/fastreports/status', { headers: headers() });
            if (!res.ok) throw new Error('status fetch failed');
            const data = await res.json();
            if (data.available) {
              el.innerHTML = '✅ <strong class="text-success">FastReport Desktop ACTIVE</strong> — DOCX/RTF/PNG/JPG are all supported.';
            } else {
              el.innerHTML = '⚠️ <strong class="text-warning">FastReport Desktop NOT INSTALLED</strong> — DOCX/RTF/PNG/JPG formats will <em>not</em> work for scheduled reports. PDF/XLSX/CSV will use legacy fallbacks.';
            }
          } catch (e) {
            el.textContent = '❌ Could not reach FastReport status endpoint.';
          }
        }
        loadTemplates().then(loadSchedules).then(loadFRStatus);
      </script>
    </body>
    </html>
    """


def _render_lookup_html() -> str:
    return """
    <!doctype html>
    <html lang="en">
    <head>
      <meta charset="utf-8"/>
      <meta name="viewport" content="width=device-width, initial-scale=1"/>
      <title>Lookup</title>
      <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
      <script src="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/js/bootstrap.bundle.min.js"></script>
      <style>
        .filter-card.active { border: 2px solid #0d6efd !important; box-shadow: 0 0.25rem 0.75rem rgba(13,110,253,.12); }
        .mono { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, "Liberation Mono", "Courier New", monospace; }
        .row-click { cursor: pointer; }
      </style>
    </head>
    <body class="bg-light">
      <nav class="navbar navbar-expand-lg navbar-dark bg-primary">
        <div class="container-fluid">
          <a class="navbar-brand" href="/dashboard">Reporting Portal</a>
          <div class="d-flex gap-2">
            <a class="btn btn-outline-light btn-sm" href="/dashboard">Dashboard</a>
            <a class="btn btn-outline-light btn-sm" href="/builder">FastReport Studio</a>
            <a class="btn btn-outline-light btn-sm" href="/schedules">Schedules</a>
          </div>
        </div>
      </nav>

      <main class="container-fluid py-3">
        <div class="row g-3">
          <div class="col-12">
            <div class="card">
              <div class="card-body">
                <div class="d-flex flex-wrap align-items-center gap-2">
                  <div class="fw-bold">Customer Lookup</div>
                  <span class="text-muted small">Search customers by phone, CNIC, name, chassis number, or engine number.</span>
                </div>
              </div>
            </div>
          </div>

          <div class="col-12 col-lg-4">
            <div id="phoneCard" class="card filter-card">
              <div class="card-header fw-bold">Phone Number Filter</div>
              <div class="card-body">
                <label class="form-label">Phone (03XXXXXXXXX)</label>
                <div class="input-group">
                  <input id="phoneInput" class="form-control mono" placeholder="03001234567" inputmode="numeric" autocomplete="off"/>
                  <button id="phoneSearchBtn" class="btn btn-primary">Search</button>
                </div>
                <div id="phoneError" class="invalid-feedback d-block"></div>
              </div>
            </div>
          </div>

          <div class="col-12 col-lg-4">
            <div id="cnicCard" class="card filter-card">
              <div class="card-header fw-bold">ID Card (CNIC) Filter</div>
              <div class="card-body">
                <label class="form-label">CNIC (12345-1234567-1)</label>
                <div class="input-group">
                  <input id="cnicInput" class="form-control mono" placeholder="12345-1234567-1" inputmode="numeric" autocomplete="off"/>
                  <button id="cnicSearchBtn" class="btn btn-primary">Search</button>
                </div>
                <div id="cnicError" class="invalid-feedback d-block"></div>
              </div>
            </div>
          </div>

          <div class="col-12 col-lg-4">
            <div id="nameCard" class="card filter-card">
              <div class="card-header fw-bold">Name Filter</div>
              <div class="card-body">
                <label class="form-label">Name (autocomplete)</label>
                <div class="input-group">
                  <input id="nameInput" class="form-control" placeholder="Type a name..." autocomplete="off" list="nameSuggestions"/>
                  <button id="nameSearchBtn" class="btn btn-primary">Search</button>
                </div>
                <datalist id="nameSuggestions"></datalist>
                <div id="nameError" class="invalid-feedback d-block"></div>
              </div>
            </div>
          </div>

          <div class="col-12 col-lg-6">
            <div id="chassisCard" class="card filter-card">
              <div class="card-header fw-bold">Chassis Number Filter</div>
              <div class="card-body">
                <label class="form-label">Chassis Number (min 3 chars, case insensitive)</label>
                <div class="input-group">
                  <input id="chassisInput" class="form-control mono" placeholder="e.g. DA232358" autocomplete="off"/>
                  <button id="chassisSearchBtn" class="btn btn-primary">Search</button>
                </div>
                <div id="chassisError" class="invalid-feedback d-block"></div>
              </div>
            </div>
          </div>

          <div class="col-12 col-lg-6">
            <div id="engineCard" class="card filter-card">
              <div class="card-header fw-bold">Engine Number Filter</div>
              <div class="card-body">
                <label class="form-label">Engine Number (min 3 chars, case insensitive)</label>
                <div class="input-group">
                  <input id="engineInput" class="form-control mono" placeholder="e.g. EN1234567" autocomplete="off"/>
                  <button id="engineSearchBtn" class="btn btn-primary">Search</button>
                </div>
                <div id="engineError" class="invalid-feedback d-block"></div>
              </div>
            </div>
          </div>

          <div class="col-12">
            <div class="card">
              <div class="card-header d-flex flex-wrap align-items-center justify-content-between gap-2">
                <div class="d-flex align-items-center gap-2">
                  <div class="fw-bold">Search Results</div>
                  <button id="printReportBtn" type="button" class="btn btn-sm btn-outline-primary">🖨️ Print Report</button>
                </div>
                <div class="d-flex align-items-center gap-2">
                  <span class="text-muted small">Filtered results</span>
                  <span id="resultsCount" class="badge text-bg-secondary">0</span>
                </div>
              </div>
              <div class="card-body p-0">
                <div class="table-responsive">
                  <table class="table table-sm table-striped align-middle mb-0" id="resultsTable">
                    <thead>
                      <tr>
                        <th role="button" class="text-nowrap" data-sort="customer_name">Customer Name</th>
                        <th role="button" class="text-nowrap" data-sort="father_name">Father's Name</th>
                        <th role="button" class="text-nowrap mono" data-sort="mobile_number">Mobile Number</th>
                        <th role="button" class="text-nowrap mono" data-sort="invoice_number">Invoice Number</th>
                        <th role="button" class="text-nowrap" data-sort="bike_model">Bike Model</th>
                        <th role="button" class="text-nowrap mono" data-sort="chassis_number">Chassis Number</th>
                        <th role="button" class="text-nowrap mono" data-sort="engine_number">Engine Number</th>
                        <th role="button" class="text-nowrap" data-sort="date">Date</th>
                      </tr>
                    </thead>
                    <tbody id="resultsBody">
                      <tr><td colspan="8" class="text-muted p-3">No results found.</td></tr>
                    </tbody>
                  </table>
                </div>
              </div>
              <div class="card-footer d-flex flex-wrap align-items-center justify-content-between gap-2">
                <div class="d-flex flex-wrap align-items-center gap-2">
                  <div class="text-muted small" id="resultsMeta"></div>
                  <span class="text-muted small">|</span>
                  <div class="text-muted small" id="selectedMeta">No row selected</div>
                  <div class="btn-group btn-group-sm ms-2" role="group" aria-label="Row actions">
                    <button id="retryUploadBtn" type="button" class="btn btn-outline-primary" disabled>Retry Upload</button>
                    <button id="copyInvoiceBtn" type="button" class="btn btn-outline-secondary" disabled>Copy Invoice #</button>
                    <button id="copyChassisBtn" type="button" class="btn btn-outline-secondary" disabled>Copy Chassis</button>
                    <button id="copyEngineBtn" type="button" class="btn btn-outline-secondary" disabled>Copy Engine</button>
                  </div>
                </div>
                <nav>
                  <ul class="pagination pagination-sm mb-0" id="pager"></ul>
                </nav>
              </div>
            </div>
          </div>
        </div>
      </main>

      <div class="modal fade" id="invoiceDetailModal" tabindex="-1" aria-labelledby="invoiceDetailTitle" aria-hidden="true">
        <div class="modal-dialog modal-lg modal-dialog-scrollable">
          <div class="modal-content">
            <div class="modal-header">
              <h5 class="modal-title" id="invoiceDetailTitle">Invoice Details</h5>
              <button type="button" class="btn-close" data-bs-dismiss="modal" aria-label="Close"></button>
            </div>
            <div class="modal-body" id="invoiceDetailBody">
              <div class="d-flex justify-content-center py-4">
                <div class="spinner-border text-primary" role="status" aria-label="Loading"></div>
              </div>
            </div>
            <div class="modal-footer">
              <button type="button" class="btn btn-outline-primary" onclick="printInvoice()">Print</button>
              <button type="button" class="btn btn-secondary" data-bs-dismiss="modal">Close</button>
            </div>
          </div>
        </div>
      </div>

      <script>
        const state = { role: 'sales', token: '', sortBy: 'date', sortDir: 'desc', page: 1, pageSize: 20, lastQueryKey: '', currentItems: [], selectedIndex: -1 };
        function headers() {
          const h = { 'Content-Type': 'application/json' };
          if (state.role) h['X-User-Role'] = state.role;
          if (state.token) h['X-API-Key'] = state.token;
          return h;
        }

        function setActive(cardId, active) {
          const el = document.getElementById(cardId);
          if (!el) return;
          if (active) el.classList.add('active');
          else el.classList.remove('active');
        }

        function setError(elId, message) {
          const el = document.getElementById(elId);
          if (!el) return;
          el.textContent = message || '';
        }

        function escapeHtml(s) {
          return String(s ?? '').replace(/[&<>"']/g, (c) => ({'&':'&amp;','<':'&lt;','>':'&gt;','"':'&quot;',"'":'&#39;'}[c]));
        }

        function getSelectedRow() {
          if (!Array.isArray(state.currentItems)) return null;
          if (state.selectedIndex < 0 || state.selectedIndex >= state.currentItems.length) return null;
          return state.currentItems[state.selectedIndex];
        }

        async function copyText(value) {
          const v = String(value || '');
          if (!v) return false;
          try {
            if (navigator.clipboard && navigator.clipboard.writeText) {
              await navigator.clipboard.writeText(v);
              return true;
            }
          } catch (e) {}
          try {
            const ta = document.createElement('textarea');
            ta.value = v;
            ta.style.position = 'fixed';
            ta.style.left = '-9999px';
            document.body.appendChild(ta);
            ta.focus();
            ta.select();
            const ok = document.execCommand('copy');
            ta.remove();
            return ok;
          } catch (e) {
            return false;
          }
        }

        function printInvoice() {
          const bodyHtml = document.getElementById('invoiceDetailBody').innerHTML;
          const title = (document.getElementById('invoiceDetailTitle').textContent || 'Invoice').trim();
          const html = `<!doctype html><html lang="en"><head><meta charset="utf-8"/>
            <title>${escapeHtml(title)}</title>
            <style>
              * { box-sizing: border-box; }
              body { font-family: "Segoe UI", Arial, Helvetica, sans-serif; color: #111; margin: 0; padding: 16mm; font-size: 11pt; }
              h2 { margin: 0 0 12px 0; font-size: 18pt; }
              .mono { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, monospace; }
              hr { margin: 1rem 0; }
              .badge { border: 1px solid #ccc; color: #111 !important; background: none !important; padding: 2px 6px; border-radius: 6px; }
              .screen-toolbar {
                position: sticky; top: 0; z-index: 10; display: flex; justify-content: space-between; align-items: center;
                padding: 10px 16px; margin: -16mm -16mm 12mm -16mm; background: #0f172a; color: white;
                border-bottom: 1px solid #ccc;
              }
              .screen-toolbar .title-text { font-weight: 600; }
              .screen-toolbar button {
                border: 0; border-radius: 6px; padding: 7px 14px; font-weight: 600; cursor: pointer; margin-left: 8px;
              }
              .btn-print { background: #2563eb; color: white; }
              .btn-close { background: #cbd5e1; color: #0f172a; }
              .hint { margin-top: 6px; font-size: 10pt; color: #fff; }
              @media print {
                body { padding: 8mm; }
                .screen-toolbar { display: none !important; }
              }
            </style></head><body>
            <div class="screen-toolbar">
              <div>
                <div class="title-text">${escapeHtml(title)}</div>
                <div class="hint">If print dialog did not open automatically, click Print below (or press Ctrl+P).</div>
              </div>
              <div>
                <button type="button" class="btn-close" onclick="window.close()">Close</button>
                <button type="button" class="btn-print" onclick="window.focus(); try { window.print(); } catch(e){}">🖨️ Print / Save as PDF</button>
              </div>
            </div>
            <div class="container py-4">
              <h2 class="mb-4">${escapeHtml(title)}</h2>
              ${bodyHtml}
            </div>
            <script>
              (function(){
                function tryPrint(){ try { window.focus(); window.print(); } catch(e){} }
                if (document.readyState === 'complete') { setTimeout(tryPrint, 300); }
                else { window.addEventListener('load', function(){ setTimeout(tryPrint, 300); }, { once: true }); }
              })();
            <\/script>
          </body></html>`;

          try {
            const blob = new Blob([html], { type: 'text/html;charset=utf-8' });
            const url = URL.createObjectURL(blob);
            const pw = window.open(url, '_blank', 'noopener');
            if (!pw) {
              URL.revokeObjectURL(url);
              alert('Please allow pop-ups for this page to open the print window.');
              return;
            }
            pw.addEventListener('pagehide', function() { try { URL.revokeObjectURL(url); } catch (e) {} });
            setTimeout(function() { try { URL.revokeObjectURL(url); } catch (e) {} }, 60000);
          } catch (e) {
            alert('Failed to open print window: ' + e);
          }
        }

        async function printReport() {
          const filters = getFilters();
          const sortBy = state.sortBy || 'date';
          const sortDir = state.sortDir || 'desc';
          const hasAnyFilter = !!(filters.phone || filters.cnic || filters.name || filters.chassis || filters.engine);
          if (!hasAnyFilter) {
            const ok = window.confirm(
              'No filters are currently set. This will print ALL invoice records in the database (up to 10,000 rows).\\n\\nDo you want to continue?'
            );
            if (!ok) return;
          }
          const btn = document.getElementById('printReportBtn');
          const origText = btn ? btn.textContent : '';
          try {
            if (btn) { btn.disabled = true; btn.textContent = 'Preparing…'; }
            const base = '/api/lookup/search?page=1&page_size=10000';
            const url =
              `${base}&phone=${encodeURIComponent(filters.phone)}&cnic=${encodeURIComponent(filters.cnic)}` +
              `&name=${encodeURIComponent(filters.name)}&chassis=${encodeURIComponent(filters.chassis)}&engine=${encodeURIComponent(filters.engine)}` +
              `&sort_by=${encodeURIComponent(sortBy)}&sort_dir=${encodeURIComponent(sortDir)}`;
            const res = await fetch(url, { headers: headers() });
            let data = null;
            try { data = await res.json(); } catch (e) { data = null; }
            if (!res.ok) {
              let detail = '';
              if (data && data.detail) {
                if (Array.isArray(data.detail)) {
                  detail = data.detail.map(d => d && d.msg ? d.msg : String(d)).join('; ');
                } else {
                  detail = String(data.detail);
                }
              }
              const statusMsg = detail ? ` (${res.status}: ${detail})` : ` (HTTP ${res.status})`;
              alert('Unable to fetch report for printing.' + statusMsg);
              return;
            }
            const payload = data || {};
            const count = payload.count || 0;
            const items = Array.isArray(payload.items) ? payload.items : [];

            const filterRows = [];
            if (filters.phone) filterRows.push(`<tr><td class="pe-3 fw-semibold">Phone</td><td>${escapeHtml(filters.phone)}</td></tr>`);
            if (filters.cnic) filterRows.push(`<tr><td class="pe-3 fw-semibold">CNIC</td><td>${escapeHtml(filters.cnic)}</td></tr>`);
            if (filters.name) filterRows.push(`<tr><td class="pe-3 fw-semibold">Name</td><td>${escapeHtml(filters.name)}</td></tr>`);
            if (filters.chassis) filterRows.push(`<tr><td class="pe-3 fw-semibold">Chassis</td><td>${escapeHtml(filters.chassis)}</td></tr>`);
            if (filters.engine) filterRows.push(`<tr><td class="pe-3 fw-semibold">Engine</td><td>${escapeHtml(filters.engine)}</td></tr>`);
            const filtersHtml = filterRows.length
              ? `<table class="table table-sm table-borderless mb-3"><tbody>${filterRows.join('')}</tbody></table>`
              : `<div class="text-muted mb-3 small">No filters applied — showing all invoices (up to 10,000 rows).</div>`;

            const cols = [
              { key: 'customer_name', label: 'Customer Name' },
              { key: 'father_name',   label: 'Father\\'s Name' },
              { key: 'mobile_number', label: 'Mobile Number', mono: true },
              { key: 'invoice_number', label: 'Invoice Number', mono: true },
              { key: 'bike_model',    label: 'Bike Model' },
              { key: 'chassis_number', label: 'Chassis Number', mono: true },
              { key: 'engine_number', label: 'Engine Number', mono: true },
              { key: 'date',          label: 'Date' },
            ];
            const headerHtml = `<tr>${cols.map(c => `<th class="text-start border-bottom border-2 border-black">${c.label}</th>`).join('')}</tr>`;
            let bodyHtml = '';
            if (!items.length) {
              bodyHtml = `<tr><td colspan="${cols.length}" class="text-muted py-4 text-center">No results found.</td></tr>`;
            } else {
              bodyHtml = items.map(r => {
                return '<tr class="align-top">' + cols.map(c => {
                  const v = (r && typeof r === 'object') ? (r[c.key] ?? '') : '';
                  const cls = c.mono ? ' class="mono"' : '';
                  return `<td${cls}>${escapeHtml(String(v))}</td>`;
                }).join('') + '</tr>';
              }).join('');
            }
            const totalAmount = items.reduce((s, r) => s + (Number(r.sale_value) || 0), 0);
            const nowStr = new Date().toLocaleString();
            const sortLabel = cols.find(c => c.key === sortBy) ? cols.find(c => c.key === sortBy).label : 'Date';

            const printTitle = 'Customer Lookup Report';
            const html = `<!doctype html><html lang="en"><head><meta charset="utf-8"/>
              <title>${printTitle}</title>
              <style>
                * { box-sizing: border-box; }
                body { font-family: "Segoe UI", Arial, Helvetica, sans-serif; color: #111; margin: 0; padding: 16mm; font-size: 11pt; }
                h1 { margin: 0 0 4mm 0; font-size: 20pt; }
                .meta { font-size: 9pt; color: #444; margin-bottom: 6mm; line-height: 1.5; }
                .mono { font-family: ui-monospace, SFMono-Regular, Menlo, Monaco, Consolas, monospace; }
                table { width: 100%; border-collapse: collapse; table-layout: fixed; }
                th, td { padding: 2.2mm 2.5mm; border-bottom: 1px solid #d0d0d0; vertical-align: top; word-wrap: break-word; text-align: left; }
                thead th { background: #0f172a; color: white; font-weight: 600; border-bottom: 1.5px solid #000; }
                tbody tr:nth-child(even) td { background: #f7f7f7; }
                .sum-table { margin-top: 5mm; width: auto; }
                .sum-table td { padding: 1.5mm 3mm; border: 0; }
                .sum-table td.k { font-weight: 600; text-align: right; }
                .card { border: 1px solid #dedede; background: #f7f7f7; padding: 10px 14px; border-radius: 8px; margin-bottom: 10px; }
                .fw-semibold { font-weight: 600; }
                .mb-2 { margin-bottom: 8px; }
                .screen-toolbar {
                  position: sticky; top: 0; z-index: 10; display: flex; justify-content: space-between; align-items: center;
                  padding: 10px 16px; margin: -16mm -16mm 12mm -16mm; background: #0f172a; color: white;
                  border-bottom: 1px solid #ccc;
                }
                .screen-toolbar .title-text { font-weight: 600; }
                .screen-toolbar button {
                  border: 0; border-radius: 6px; padding: 7px 14px; font-weight: 600; cursor: pointer; margin-left: 8px;
                }
                .btn-print { background: #2563eb; color: white; }
                .btn-close { background: #cbd5e1; color: #0f172a; }
                .hint { margin-top: 6px; font-size: 10pt; color: #fff; }
                @media print {
                  body { padding: 8mm; }
                  thead { display: table-header-group; }
                  tr    { page-break-inside: avoid; }
                  .screen-toolbar { display: none !important; }
                }
              </style></head><body>
              <div class="screen-toolbar">
                <div>
                  <div class="title-text">${printTitle}</div>
                  <div class="hint">If print dialog did not open automatically, click Print below (or press Ctrl+P).</div>
                </div>
                <div>
                  <button type="button" class="btn-close" onclick="window.close()">Close</button>
                  <button type="button" class="btn-print" onclick="window.focus(); try { window.print(); } catch(e){}">🖨️ Print / Save as PDF</button>
                </div>
              </div>
              <h1>${printTitle}</h1>
              <div class="meta">
                <div>Generated: <span class="mono">${escapeHtml(nowStr)}</span></div>
                <div>Total records: <b>${count}</b> &nbsp;|&nbsp; Sorted by: <b>${escapeHtml(sortLabel)} (${escapeHtml(sortDir.toUpperCase())})</b></div>
              </div>
              <div class="card">
                <div class="fw-semibold mb-2">Applied Filters</div>
                ${filtersHtml}
              </div>
              <table>
                <thead>${headerHtml}</thead>
                <tbody>${bodyHtml}</tbody>
              </table>
              <table class="sum-table">
                <tr><td class="k">Total rows printed:</td><td>${count}</td></tr>
                <tr><td class="k">Aggregate sale value:</td><td class="mono">Rs ${totalAmount.toLocaleString(undefined, { minimumFractionDigits: 2, maximumFractionDigits: 2 })}</td></tr>
              </table>
              <script>
                (function(){
                  function tryPrint(){ try { window.focus(); window.print(); } catch(e){} }
                  if (document.readyState === 'complete') { setTimeout(tryPrint, 450); }
                  else { window.addEventListener('load', function(){ setTimeout(tryPrint, 450); }, { once: true }); }
                })();
              <\/script>
            </body></html>`;

            try {
              const blob = new Blob([html], { type: 'text/html;charset=utf-8' });
              const url = URL.createObjectURL(blob);
              const pw = window.open(url, '_blank', 'noopener');
              if (!pw) {
                URL.revokeObjectURL(url);
                alert('Please allow pop-ups for this page to open the print window.');
                return;
              }
              pw.addEventListener('pagehide', function() { try { URL.revokeObjectURL(url); } catch (e) {} });
              setTimeout(function() { try { URL.revokeObjectURL(url); } catch (e) {} }, 60000);
            } catch (e) {
              alert(`Print failed: ${e}`);
            }
          } finally {
            if (btn) { btn.disabled = false; btn.textContent = origText || '🖨️ Print Report'; }
          }
        }

        function updateActionButtons() {
          const row = getSelectedRow();
          const selectedMeta = document.getElementById('selectedMeta');
          const retryBtn = document.getElementById('retryUploadBtn');
          const copyInvBtn = document.getElementById('copyInvoiceBtn');
          const copyChBtn = document.getElementById('copyChassisBtn');
          const copyEnBtn = document.getElementById('copyEngineBtn');

          if (!row) {
            selectedMeta.textContent = 'No row selected';
            retryBtn.disabled = true;
            copyInvBtn.disabled = true;
            copyChBtn.disabled = true;
            if (copyEnBtn) copyEnBtn.disabled = true;
            return;
          }

          const inv = row.invoice_number || '';
          const ch = row.chassis_number || '';
          const en = row.engine_number || '';
          const st = String(row.sync_status || '').toUpperCase();
          selectedMeta.textContent = inv ? `Selected: ${inv}` : 'Selected row';
          copyInvBtn.disabled = !inv;
          copyChBtn.disabled = !ch;
          if (copyEnBtn) copyEnBtn.disabled = !en;
          retryBtn.disabled = !inv || !['PENDING', 'FAILED'].includes(st);
        }

        async function retrySelected() {
          const row = getSelectedRow();
          if (!row) return;
          const inv = row.invoice_number || '';
          if (!inv) return;

          const btn = document.getElementById('retryUploadBtn');
          btn.disabled = true;
          btn.textContent = 'Retrying...';
          try {
            const res = await fetch(`/api/invoices/${encodeURIComponent(inv)}/retry`, { method: 'POST', headers: headers() });
            const data = await res.json();
            if (!res.ok) {
              setError('nameError', data && data.detail ? data.detail : 'Retry failed.');
              return;
            }
            await runSearch(false);
          } catch (e) {
            setError('nameError', `Network error: ${e}`);
          } finally {
            btn.textContent = 'Retry Upload';
            updateActionButtons();
          }
        }

        function showInvoiceModal(titleText) {
          const modalEl = document.getElementById('invoiceDetailModal');
          const titleEl = document.getElementById('invoiceDetailTitle');
          if (titleEl) titleEl.textContent = titleText || 'Invoice Details';
          const modal = bootstrap.Modal.getOrCreateInstance(modalEl, { backdrop: true, focus: true, keyboard: true });
          modal.show();
          return modal;
        }

        function setInvoiceModalBody(html) {
          const body = document.getElementById('invoiceDetailBody');
          body.innerHTML = html;
        }

        function renderInvoiceDetails(inv) {
          const customer = inv.customer || {};
          const items = Array.isArray(inv.items) ? inv.items : [];
          const totals = inv.totals || {};
          const status = String(inv.sync_status || '').toUpperCase();
          const statusClass =
            status === 'SUCCESS' ? 'text-bg-success' :
            status === 'PENDING' ? 'text-bg-warning' :
            status === 'FAILED' ? 'text-bg-danger' : 'text-bg-secondary';

          const itemsRows = items.map((it) => `
            <tr>
              <td>${escapeHtml(it.item_name || it.item_code || '')}</td>
              <td class="mono">${escapeHtml(it.item_code || '')}</td>
              <td class="text-end">${escapeHtml(it.quantity ?? '')}</td>
              <td class="text-end">${escapeHtml((it.sale_value ?? '').toString())}</td>
              <td class="text-end">${escapeHtml((it.tax_charged ?? '').toString())}</td>
              <td class="text-end">${escapeHtml((it.further_tax ?? '').toString())}</td>
              <td class="text-end">${escapeHtml((it.total_amount ?? '').toString())}</td>
              <td>${escapeHtml(it.model || '')}</td>
              <td class="mono">${escapeHtml(it.chassis_number || '')}</td>
              <td class="mono">${escapeHtml(it.engine_number || '')}</td>
            </tr>
          `).join('');

          const fbrInfo = inv.fbr_invoice_number ? `<div class="text-muted small">FBR Invoice: <span class="mono">${escapeHtml(inv.fbr_invoice_number)}</span></div>` : '';

          return `
            <div class="d-flex flex-wrap align-items-center justify-content-between gap-2">
              <div>
                <div class="fw-bold">${escapeHtml(inv.invoice_number || '')}</div>
                <div class="text-muted small">${escapeHtml(inv.date || '')}</div>
                ${fbrInfo}
              </div>
              <div class="d-flex align-items-center gap-2">
                <span class="badge ${statusClass}">${escapeHtml(status || 'UNKNOWN')}</span>
              </div>
            </div>

            <hr/>

            <div class="row g-3">
              <div class="col-12 col-lg-6">
                <div class="fw-bold mb-2">Customer</div>
                <div class="small">
                  <div><span class="text-muted">Name:</span> ${escapeHtml(customer.name || '')}</div>
                  <div><span class="text-muted">Father:</span> ${escapeHtml(customer.father_name || '')}</div>
                  <div><span class="text-muted">CNIC:</span> <span class="mono">${escapeHtml(customer.cnic || '')}</span></div>
                  <div><span class="text-muted">NTN:</span> <span class="mono">${escapeHtml(customer.ntn || '')}</span></div>
                  <div><span class="text-muted">Phone:</span> <span class="mono">${escapeHtml(customer.phone || '')}</span></div>
                  <div><span class="text-muted">Type:</span> ${escapeHtml(customer.type || '')}</div>
                  <div><span class="text-muted">Address:</span> ${escapeHtml(customer.address || '')}</div>
                </div>
              </div>
              <div class="col-12 col-lg-6">
                <div class="fw-bold mb-2">Invoice</div>
                <div class="small">
                  <div><span class="text-muted">POS ID:</span> <span class="mono">${escapeHtml(inv.pos_id || '')}</span></div>
                  <div><span class="text-muted">Payment:</span> ${escapeHtml(inv.payment_mode || '')}</div>
                  <div><span class="text-muted">Fiscalized:</span> ${inv.is_fiscalized ? 'Yes' : 'No'}</div>
                </div>
              </div>
            </div>

            <hr/>

            <div class="fw-bold mb-2">Items</div>
            <div class="table-responsive">
              <table class="table table-sm table-striped align-middle mb-0">
                <thead>
                  <tr>
                    <th>Item</th>
                    <th class="mono">Code</th>
                    <th class="text-end">Qty</th>
                    <th class="text-end">Subtotal</th>
                    <th class="text-end">Tax</th>
                    <th class="text-end">Further Tax</th>
                    <th class="text-end">Total</th>
                    <th>Model</th>
                    <th class="mono">Chassis</th>
                    <th class="mono">Engine</th>
                  </tr>
                </thead>
                <tbody>
                  ${itemsRows || `<tr><td colspan="10" class="text-muted">No items.</td></tr>`}
                </tbody>
              </table>
            </div>

            <div class="d-flex justify-content-end mt-3">
              <table class="table table-sm w-auto mb-0">
                <tbody>
                  <tr><td class="text-muted">Subtotal</td><td class="text-end">${escapeHtml((totals.subtotal ?? '').toString())}</td></tr>
                  <tr><td class="text-muted">Tax</td><td class="text-end">${escapeHtml((totals.tax ?? '').toString())}</td></tr>
                  <tr><td class="text-muted">Further Tax</td><td class="text-end">${escapeHtml((totals.further_tax ?? '').toString())}</td></tr>
                  <tr class="fw-bold"><td>Total</td><td class="text-end">${escapeHtml((totals.total ?? '').toString())}</td></tr>
                </tbody>
              </table>
            </div>
          `;
        }

        async function openInvoiceDetails(invoiceNumber) {
          const inv = String(invoiceNumber || '').trim();
          if (!inv) return;
          showInvoiceModal(`Invoice Details - ${inv}`);
          setInvoiceModalBody(`
            <div class="d-flex justify-content-center py-4">
              <div class="spinner-border text-primary" role="status" aria-label="Loading invoice details"></div>
            </div>
          `);
          try {
            const res = await fetch(`/api/invoices/${encodeURIComponent(inv)}/details`, { headers: headers() });
            const data = await res.json();
            if (!res.ok) {
              const msg = data && data.detail ? data.detail : 'Unable to load invoice details.';
              setInvoiceModalBody(`<div class="alert alert-danger" role="alert">${escapeHtml(msg)}</div>`);
              return;
            }
            setInvoiceModalBody(renderInvoiceDetails(data));
          } catch (e) {
            setInvoiceModalBody(`<div class="alert alert-danger" role="alert">Network error: ${escapeHtml(e)}</div>`);
          }
        }

        function normalizePhone(raw) {
          const digits = String(raw || '').replace(/\\D/g, '').slice(0, 11);
          return digits;
        }

        function normalizeCnic(raw) {
          const digits = String(raw || '').replace(/\\D/g, '').slice(0, 13);
          let out = digits;
          if (digits.length > 5) out = digits.slice(0,5) + '-' + digits.slice(5);
          if (digits.length > 12) out = out.slice(0,13) + '-' + out.slice(13);
          return out;
        }

        function validPhone(phone) {
          return /^03\\d{9}$/.test(phone);
        }

        function validCnic(cnic) {
          return /^\\d{5}-\\d{7}-\\d$/.test(cnic);
        }

        function normalizeChassis(raw) {
          return String(raw || '').replace(/[^A-Za-z0-9]/g, '').toUpperCase().slice(0, 50);
        }

        function normalizeEngine(raw) {
          return String(raw || '').replace(/[^A-Za-z0-9]/g, '').toUpperCase().slice(0, 50);
        }

        let nameTimer = null;
        async function autocompleteName() {
          const q = (document.getElementById('nameInput').value || '').trim();
          if (q.length < 2) return;
          const res = await fetch(`/api/customers/autocomplete?query=${encodeURIComponent(q)}`, { headers: headers() });
          const data = await res.json();
          if (!res.ok) return;
          const list = document.getElementById('nameSuggestions');
          list.innerHTML = '';
          (data.items || []).forEach((v) => {
            const opt = document.createElement('option');
            opt.value = v;
            list.appendChild(opt);
          });
        }

        function getFilters() {
          const phone = normalizePhone(document.getElementById('phoneInput').value || '');
          const cnic = normalizeCnic(document.getElementById('cnicInput').value || '');
          const name = (document.getElementById('nameInput').value || '').trim();
          const chassis = normalizeChassis(document.getElementById('chassisInput') ? document.getElementById('chassisInput').value : '');
          const engine = normalizeEngine(document.getElementById('engineInput') ? document.getElementById('engineInput').value : '');
          return { phone, cnic, name, chassis, engine };
        }

        function validateFilters(filters) {
          setError('phoneError', '');
          setError('cnicError', '');
          setError('nameError', '');
          setError('chassisError', '');
          setError('engineError', '');

          let ok = true;
          if (filters.phone && !validPhone(filters.phone)) {
            setError('phoneError', 'Invalid phone number. Use 03XXXXXXXXX (11 digits).');
            ok = false;
          }
          if (filters.cnic && !validCnic(filters.cnic)) {
            setError('cnicError', 'Invalid CNIC format. Use 12345-1234567-1.');
            ok = false;
          }
          if (filters.name && filters.name.length < 2) {
            setError('nameError', 'Please type at least 2 characters.');
            ok = false;
          }
          if (filters.chassis && filters.chassis.length < 3) {
            setError('chassisError', 'Chassis Number must be at least 3 characters.');
            ok = false;
          }
          if (filters.engine && filters.engine.length < 3) {
            setError('engineError', 'Engine Number must be at least 3 characters.');
            ok = false;
          }

          setActive('phoneCard', !!filters.phone && validPhone(filters.phone));
          setActive('cnicCard', !!filters.cnic && validCnic(filters.cnic));
          setActive('nameCard', !!filters.name && filters.name.length >= 2);
          setActive('chassisCard', !!filters.chassis && filters.chassis.length >= 3);
          setActive('engineCard', !!filters.engine && filters.engine.length >= 3);

          const hasAny = !!(filters.phone || filters.cnic || filters.name || filters.chassis || filters.engine);
          return hasAny && ok;
        }

        function setLoading(isLoading) {
          const meta = document.getElementById('resultsMeta');
          if (!meta) return;
          meta.textContent = isLoading ? 'Loading…' : '';
        }

        function renderTable(payload) {
          const body = document.getElementById('resultsBody');
          const count = payload && typeof payload.count === 'number' ? payload.count : 0;
          const items = payload && Array.isArray(payload.items) ? payload.items : [];
          state.currentItems = items;
          state.selectedIndex = -1;
          document.getElementById('resultsCount').textContent = String(count);

          if (!items.length) {
            body.innerHTML = `<tr><td colspan="8" class="text-muted p-3">No results found.</td></tr>`;
          } else {
            const rows = items.map((r) => `
              <tr class="row-click" role="button" tabindex="0" aria-label="View invoice ${escapeHtml(r.invoice_number || '')}">
                <td>${escapeHtml(r.customer_name || '')}</td>
                <td>${escapeHtml(r.father_name || '')}</td>
                <td class="mono">${escapeHtml(r.mobile_number || '')}</td>
                <td class="mono">${escapeHtml(r.invoice_number || '')}</td>
                <td>${escapeHtml(r.bike_model || '')}</td>
                <td class="mono">${escapeHtml(r.chassis_number || '')}</td>
                <td class="mono">${escapeHtml(r.engine_number || '')}</td>
                <td>${escapeHtml(r.date || '')}</td>
              </tr>
            `).join('');
            body.innerHTML = rows;
          }
          document.querySelectorAll('#resultsBody tr').forEach((tr, idx) => {
            if (!state.currentItems.length) return;
            const open = () => openInvoiceDetails(state.currentItems[idx].invoice_number || '');
            tr.addEventListener('click', () => {
              state.selectedIndex = idx;
              document.querySelectorAll('#resultsBody tr').forEach((x) => x.classList.remove('table-active'));
              tr.classList.add('table-active');
              updateActionButtons();
              open();
            });
            tr.addEventListener('keydown', (e) => {
              if (e.key === 'Enter' || e.key === ' ') {
                e.preventDefault();
                state.selectedIndex = idx;
                document.querySelectorAll('#resultsBody tr').forEach((x) => x.classList.remove('table-active'));
                tr.classList.add('table-active');
                updateActionButtons();
                open();
              }
            });
          });
          updateActionButtons();

          const meta = document.getElementById('resultsMeta');
          const page = payload && payload.page ? payload.page : 1;
          const pageSize = payload && payload.page_size ? payload.page_size : state.pageSize;
          const totalPages = payload && payload.total_pages ? payload.total_pages : 1;
          const shownStart = count ? ((page - 1) * pageSize + 1) : 0;
          const shownEnd = Math.min(page * pageSize, count);
          meta.textContent = count ? `Showing ${shownStart}-${shownEnd} of ${count}` : 'Showing 0 results';

          const pager = document.getElementById('pager');
          pager.innerHTML = '';
          if (totalPages <= 1) return;

          function pageItem(label, targetPage, disabled, active) {
            const li = document.createElement('li');
            li.className = `page-item ${disabled ? 'disabled' : ''} ${active ? 'active' : ''}`;
            const a = document.createElement('a');
            a.className = 'page-link';
            a.href = '#';
            a.textContent = label;
            a.addEventListener('click', (e) => {
              e.preventDefault();
              if (disabled) return;
              state.page = targetPage;
              runSearch(false);
            });
            li.appendChild(a);
            return li;
          }

          pager.appendChild(pageItem('Prev', Math.max(1, page - 1), page <= 1, false));

          const start = Math.max(1, page - 2);
          const end = Math.min(totalPages, start + 4);
          for (let p = start; p <= end; p++) {
            pager.appendChild(pageItem(String(p), p, false, p === page));
          }

          pager.appendChild(pageItem('Next', Math.min(totalPages, page + 1), page >= totalPages, false));
        }

        async function runSearch(resetPage) {
          const filters = getFilters();
          document.getElementById('phoneInput').value = filters.phone;
          document.getElementById('cnicInput').value = filters.cnic;
          const chInput = document.getElementById('chassisInput');
          const enInput = document.getElementById('engineInput');
          if (chInput) chInput.value = filters.chassis;
          if (enInput) enInput.value = filters.engine;

          if (!validateFilters(filters)) {
            renderTable({ count: 0, items: [], page: 1, page_size: state.pageSize, total_pages: 1 });
            return;
          }

          if (resetPage) state.page = 1;
          const queryKey = JSON.stringify({ ...filters, sortBy: state.sortBy, sortDir: state.sortDir, page: state.page });
          if (queryKey === state.lastQueryKey) return;
          state.lastQueryKey = queryKey;
          setLoading(true);
          try {
            const url =
              `/api/lookup/search?phone=${encodeURIComponent(filters.phone)}&cnic=${encodeURIComponent(filters.cnic)}` +
              `&name=${encodeURIComponent(filters.name)}&chassis=${encodeURIComponent(filters.chassis)}&engine=${encodeURIComponent(filters.engine)}` +
              `&sort_by=${encodeURIComponent(state.sortBy)}&sort_dir=${encodeURIComponent(state.sortDir)}` +
              `&page=${encodeURIComponent(state.page)}&page_size=${encodeURIComponent(state.pageSize)}`;
            const res = await fetch(url, { headers: headers() });
            const data = await res.json();
            if (!res.ok) {
              const msg = data && data.detail ? data.detail : 'Search failed.';
              setError('nameError', msg);
              renderTable({ count: 0, items: [], page: 1, page_size: state.pageSize, total_pages: 1 });
              return;
            }
            renderTable(data);
          } catch (e) {
            setError('nameError', `Network error: ${e}`);
            renderTable({ count: 0, items: [], page: 1, page_size: state.pageSize, total_pages: 1 });
          } finally {
            setLoading(false);
          }
        }

        document.getElementById('phoneInput').addEventListener('input', (e) => {
          e.target.value = normalizePhone(e.target.value);
          setError('phoneError', '');
          if (nameTimer) clearTimeout(nameTimer);
          nameTimer = setTimeout(() => runSearch(true), 220);
        });
        document.getElementById('cnicInput').addEventListener('input', (e) => {
          e.target.value = normalizeCnic(e.target.value);
          setError('cnicError', '');
          if (nameTimer) clearTimeout(nameTimer);
          nameTimer = setTimeout(() => runSearch(true), 220);
        });
        document.getElementById('nameInput').addEventListener('input', () => {
          setError('nameError', '');
          if (nameTimer) clearTimeout(nameTimer);
          nameTimer = setTimeout(() => { autocompleteName(); runSearch(true); }, 220);
        });
        const chassisInputEl = document.getElementById('chassisInput');
        if (chassisInputEl) {
          chassisInputEl.addEventListener('input', (e) => {
            e.target.value = normalizeChassis(e.target.value);
            setError('chassisError', '');
            if (nameTimer) clearTimeout(nameTimer);
            nameTimer = setTimeout(() => runSearch(true), 220);
          });
        }
        const engineInputEl = document.getElementById('engineInput');
        if (engineInputEl) {
          engineInputEl.addEventListener('input', (e) => {
            e.target.value = normalizeEngine(e.target.value);
            setError('engineError', '');
            if (nameTimer) clearTimeout(nameTimer);
            nameTimer = setTimeout(() => runSearch(true), 220);
          });
        }

        document.getElementById('phoneSearchBtn').addEventListener('click', () => runSearch(true));
        document.getElementById('cnicSearchBtn').addEventListener('click', () => runSearch(true));
        document.getElementById('nameSearchBtn').addEventListener('click', () => runSearch(true));
        const chassisSearchBtn = document.getElementById('chassisSearchBtn');
        if (chassisSearchBtn) chassisSearchBtn.addEventListener('click', () => runSearch(true));
        const engineSearchBtn = document.getElementById('engineSearchBtn');
        if (engineSearchBtn) engineSearchBtn.addEventListener('click', () => runSearch(true));
        const printReportBtn = document.getElementById('printReportBtn');
        if (printReportBtn) printReportBtn.addEventListener('click', () => printReport());
        document.getElementById('retryUploadBtn').addEventListener('click', retrySelected);
        document.getElementById('copyInvoiceBtn').addEventListener('click', async () => {
          const row = getSelectedRow();
          if (!row) return;
          const ok = await copyText(row.invoice_number || '');
          if (!ok) setError('nameError', 'Unable to copy invoice number.');
        });
        document.getElementById('copyChassisBtn').addEventListener('click', async () => {
          const row = getSelectedRow();
          if (!row) return;
          const ok = await copyText(row.chassis_number || '');
          if (!ok) setError('nameError', 'Unable to copy chassis number.');
        });
        const copyEngineBtn = document.getElementById('copyEngineBtn');
        if (copyEngineBtn) {
          copyEngineBtn.addEventListener('click', async () => {
            const row = getSelectedRow();
            if (!row) return;
            const ok = await copyText(row.engine_number || '');
            if (!ok) setError('nameError', 'Unable to copy engine number.');
          });
        }

        document.querySelectorAll('#resultsTable thead th[data-sort]').forEach((th) => {
          th.addEventListener('click', () => {
            const key = th.getAttribute('data-sort');
            if (!key) return;
            if (state.sortBy === key) state.sortDir = state.sortDir === 'asc' ? 'desc' : 'asc';
            else { state.sortBy = key; state.sortDir = 'asc'; }
            state.page = 1;
            runSearch(false);
          });
        });

        runSearch(true);
      </script>
    </body>
    </html>
    """


@app.on_event("startup")
def _startup() -> None:
    if (os.getenv("REPORTING_DISABLE_SCHEDULER") or "").strip() == "1":
        return

    root = _reporting_root_dir()
    root.mkdir(parents=True, exist_ok=True)

    def loop() -> None:
        while True:
            try:
                _run_due_schedules()
            except Exception:
                pass
            time.sleep(30)

    threading.Thread(target=loop, daemon=True, name="ReportingScheduler").start()


def _run_due_schedules() -> None:
    db = SessionLocal()
    try:
        _load_or_create_default_template(db)
        schedules = db.query(ReportSchedule).filter(ReportSchedule.enabled.is_(True)).all()
        now = datetime.utcnow()
        for sch in schedules:
            last = sch.last_run_at
            interval = int(sch.interval_minutes or 60)
            due = (not last) or ((now - last).total_seconds() >= interval * 60)
            if not due:
                continue
            sch.last_run_at = now
            db.add(sch)
            db.commit()
            _execute_schedule(db, sch)
    finally:
        db.close()


def _smtp_config() -> Dict[str, Any]:
    return {
        "host": os.getenv("REPORTING_SMTP_HOST", ""),
        "port": int(os.getenv("REPORTING_SMTP_PORT", "587") or "587"),
        "user": os.getenv("REPORTING_SMTP_USER", ""),
        "password": os.getenv("REPORTING_SMTP_PASSWORD", ""),
        "from": os.getenv("REPORTING_SMTP_FROM", ""),
        "use_tls": (os.getenv("REPORTING_SMTP_TLS", "1") or "1") == "1",
    }


def _send_email(to_list: List[str], subject: str, body: str, attachment_name: str, attachment_bytes: bytes) -> None:
    cfg = _smtp_config()
    if not cfg["host"] or not cfg["from"] or not to_list:
        return

    msg = EmailMessage()
    msg["From"] = cfg["from"]
    msg["To"] = ", ".join(to_list)
    msg["Subject"] = subject
    msg.set_content(body)
    msg.add_attachment(attachment_bytes, maintype="application", subtype="octet-stream", filename=attachment_name)

    if cfg["use_tls"]:
        with smtplib.SMTP(cfg["host"], cfg["port"]) as s:
            s.starttls()
            if cfg["user"]:
                s.login(cfg["user"], cfg["password"])
            s.send_message(msg)
    else:
        with smtplib.SMTP(cfg["host"], cfg["port"]) as s:
            if cfg["user"]:
                s.login(cfg["user"], cfg["password"])
            s.send_message(msg)


def _execute_schedule(db: Session, sch: ReportSchedule) -> None:
    run = ReportRun(schedule_id=sch.id, status="STARTED")
    db.add(run)
    db.commit()
    db.refresh(run)

    try:
        tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == sch.template_id).first()
        tmpl = _ensure_template_has_widgets(db, tmpl)

        metrics = _compute_metrics(db, None, None, "ALL")
        file_bytes, name, media_type = _export_bytes(tmpl, metrics, sch.export_format)

        root = _reporting_root_dir()
        root.mkdir(parents=True, exist_ok=True)
        file_path = root / f"schedule_{sch.id}_run_{run.id}_{name}"
        file_path.write_bytes(file_bytes)

        run.status = "SUCCESS"
        run.finished_at = datetime.utcnow()
        run.file_path = str(file_path)
        db.add(run)
        db.commit()

        recipients = sch.recipients or []
        if isinstance(recipients, list) and recipients:
            _send_email(
                [str(x) for x in recipients],
                subject=f"Scheduled Report: {tmpl.name}",
                body=f"Scheduled report generated at {datetime.utcnow().isoformat(sep=' ')}",
                attachment_name=name,
                attachment_bytes=file_bytes,
            )

    except Exception as e:
        run.status = "FAILED"
        run.finished_at = datetime.utcnow()
        run.error_message = str(e)
        db.add(run)
        db.commit()


def _export_csv(metrics: Dict[str, Any]) -> bytes:
    buffer = StringIO()
    buffer.write("invoice_number,usin,fbr_invoice_number,datetime,pos_id,payment_mode,total_amount,sync_status\n")
    for r in metrics.get("invoices", []):
        buffer.write(
            f"{r['invoice_number']},{r.get('usin','')},{r.get('fbr_invoice_number','')},{r['datetime']},{r['pos_id']},{r['payment_mode']},{float(r['total_amount']):.2f},{r['sync_status']}\n"
        )
    return buffer.getvalue().encode("utf-8")


def _export_xlsx(metrics: Dict[str, Any]) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Invoices"
    ws.append(["Invoice #", "USIN", "FBR Invoice #", "DateTime", "POS ID", "Payment Mode", "Total Amount", "Sync Status"])
    for r in metrics.get("invoices", []):
        ws.append([r["invoice_number"], r.get("usin",""), r.get("fbr_invoice_number",""), r["datetime"], r["pos_id"], r["payment_mode"], float(r["total_amount"]), r["sync_status"]])
    bio = BytesIO()
    wb.save(bio)
    return bio.getvalue()


def _export_pdf(template: ReportTemplate, metrics: Dict[str, Any]) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas

    bio = BytesIO()
    c = canvas.Canvas(bio, pagesize=A4)
    width, height = A4
    y = height - 40
    c.setFont("Helvetica-Bold", 14)
    c.drawString(40, y, f"Report: {template.name}")
    y -= 24
    c.setFont("Helvetica", 10)
    c.drawString(40, y, f"Generated: {datetime.utcnow().isoformat(sep=' ')}")
    y -= 24
    c.setFont("Helvetica", 11)
    c.drawString(40, y, f"Total Invoices: {metrics.get('total_invoices', 0)}")
    y -= 18
    c.drawString(40, y, f"Total Amount: {float(metrics.get('total_amount', 0.0)):.2f}")
    y -= 30
    c.setFont("Helvetica-Bold", 11)
    c.drawString(40, y, "Recent Invoices")
    y -= 18
    c.setFont("Helvetica", 9)
    for r in metrics.get("invoices", [])[:25]:
        c.drawString(40, y, f"{r['invoice_number']} | {r['datetime']} | {r['total_amount']:.2f} | {r['sync_status']}")
        y -= 12
        if y < 60:
            c.showPage()
            y = height - 40
            c.setFont("Helvetica", 9)
    c.showPage()
    c.save()
    return bio.getvalue()


def _export_pptx(template: ReportTemplate, metrics: Dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = f"{template.name}"
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(2.0))
    tf = tx.text_frame
    tf.text = f"Generated: {datetime.utcnow().isoformat(sep=' ')}"
    p = tf.add_paragraph()
    p.text = f"Total Invoices: {metrics.get('total_invoices', 0)}"
    p = tf.add_paragraph()
    p.text = f"Total Amount: {float(metrics.get('total_amount', 0.0)):.2f}"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    title2 = slide2.shapes.title
    if title2:
        title2.text = "Recent Invoices"
    box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.5), Inches(5.0))
    tf2 = box.text_frame
    tf2.word_wrap = True
    for r in metrics.get("invoices", [])[:15]:
        p = tf2.add_paragraph()
        p.text = f"{r['invoice_number']} | {r['datetime']} | {r['total_amount']:.2f} | {r['sync_status']}"

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def _frx_template_for_format(fmt: str, template_name_hint: str = "") -> str:
    """Pick the best .frx template name for a given export format/context."""
    name = (template_name_hint or "").lower()
    if "ledger" in name:
        return "customer_ledger"
    if "invoice" in name or "bill" in name:
        return "invoice"
    if "authority" in name:
        return "authority_letter"
    if "lookup" in name or "customer" in name:
        return "customer_lookup_report"
    return "sales_dashboard"


def _try_fastreports_export(
    fmt: str,
    metrics: Dict[str, Any],
    template: Optional[ReportTemplate] = None,
) -> Optional[Tuple[bytes, str, str]]:
    """Attempt FastReport Desktop export. Returns None if unavailable/fails."""
    if _fr_available is None or not _fr_available():
        return None

    template_name = _frx_template_for_format(fmt, template.name if template else "")
    report_data: Dict[str, Any] = {
        "report_title": template.name if template else "Report",
        "generated_at": datetime.utcnow().isoformat(sep=" "),
        "total_invoices": metrics.get("total_invoices", 0),
        "total_amount": float(metrics.get("total_amount", 0.0)),
        "avg_invoice_amount": float(metrics.get("avg_invoice_amount", 0.0)),
        "daily_sales": metrics.get("daily_sales", []),
        "status_breakdown": metrics.get("status_breakdown", []),
        "invoices": metrics.get("invoices", []),
    }

    suffix_map = {
        "pdf": ".pdf",
        "html": ".html",
        "htm": ".html",
        "xlsx": ".xlsx",
        "excel": ".xlsx",
        "csv": ".csv",
        "docx": ".docx",
        "rtf": ".rtf",
        "pptx": ".pptx",
        "png": ".png",
        "jpg": ".jpg",
    }
    media_type_map = {
        "pdf": "application/pdf",
        "html": "text/html",
        "htm": "text/html",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "csv": "text/csv",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "rtf": "application/rtf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "png": "image/png",
        "jpg": "image/jpeg",
    }

    fmt_lower = (fmt or "pdf").lower()
    suffix = suffix_map.get(fmt_lower, f".{fmt_lower}")
    media_type = media_type_map.get(fmt_lower, "application/octet-stream")

    try:
        result = _fr_build_report(
            template_name_or_path=template_name,
            data=report_data,
            export_format=fmt_lower,
        )
        if not result.ok or not result.output_path:
            logger.warning("FastReport render failed: %s", result.error)
            return None
        out_path = result.output_path
        if not out_path.is_file():
            return None
        file_bytes = out_path.read_bytes()
        filename = f"{template_name}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}{suffix}"
        return (file_bytes, filename, media_type)
    except Exception as exc:
        logger.exception("FastReport export threw an exception (will fallback): %s", exc)
        return None


def _export_bytes(template: ReportTemplate, metrics: Dict[str, Any], fmt: str) -> Tuple[bytes, str, str]:
    fmt = (fmt or "").lower()
    fr_result = _try_fastreports_export(fmt, metrics, template)
    if fr_result is not None:
        return fr_result

    if fmt == "csv":
        return _export_csv(metrics), "report.csv", "text/csv"
    if fmt == "xlsx":
        return _export_xlsx(metrics), "report.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if fmt == "pptx":
        return _export_pptx(template, metrics), "report.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    return _export_pdf(template, metrics), "report.pdf", "application/pdf"


@app.get("/", response_class=HTMLResponse)
def root() -> RedirectResponse:
    return RedirectResponse(url="/dashboard")


@app.get("/dashboard", response_class=HTMLResponse)
def dashboard() -> str:
    return _render_dashboard_html()

@app.get("/favicon.ico")
def favicon() -> StreamingResponse:
    return StreamingResponse(BytesIO(b""), media_type="image/x-icon")


@app.get("/builder", response_class=HTMLResponse)
def builder() -> str:
    return _render_builder_html()


@app.get("/schedules", response_class=HTMLResponse)
def schedules_page() -> str:
    return _render_schedules_html()


@app.get("/lookup", response_class=HTMLResponse)
def lookup_page() -> str:
    return _render_lookup_html()


def _customer_to_dict(c: Customer) -> Dict[str, Any]:
    return {
        "id": int(c.id),
        "name": c.name,
        "father_name": c.father_name,
        "business_name": c.business_name,
        "cnic": c.cnic,
        "ntn": c.ntn,
        "phone": c.phone,
        "address": c.address,
        "type": c.type,
        "is_deleted": bool(getattr(c, "is_deleted", False)),
        "created_at": c.created_at.isoformat(sep=" ") if getattr(c, "created_at", None) else None,
    }


@app.get("/api/templates")
def list_templates(
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    items = db.query(ReportTemplate).order_by(ReportTemplate.name.asc()).all()
    return JSONResponse({"items": [{"id": t.id, "name": t.name, "description": t.description, "is_active": t.is_active} for t in items]})


@app.get("/api/templates/{template_id}")
def get_template(
    template_id: int,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == template_id).first()
    if not tmpl:
        raise HTTPException(status_code=404, detail="Not found")
    return JSONResponse({"id": tmpl.id, "name": tmpl.name, "description": tmpl.description, "definition": tmpl.definition})


@app.post("/api/templates")
def create_template(
    payload: Dict[str, Any],
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])
    name = (payload.get("name") or "").strip()
    if not name:
        raise HTTPException(status_code=400, detail="name required")
    definition = payload.get("definition") or {"version": 1, "widgets": []}
    tmpl = ReportTemplate(name=name, description=payload.get("description"), definition=definition, is_active=True, created_by_role=role)
    db.add(tmpl)
    db.commit()
    db.refresh(tmpl)
    _audit(db, action="CREATE", resource_type="REPORT_TEMPLATE", resource_id=tmpl.id, details={"name": tmpl.name}, request=request)
    return JSONResponse({"id": tmpl.id})


@app.put("/api/templates/{template_id}")
def update_template(
    template_id: int,
    payload: Dict[str, Any],
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])
    tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == template_id).first()
    if not tmpl:
        raise HTTPException(status_code=404, detail="Not found")
    name = (payload.get("name") or tmpl.name).strip()
    tmpl.name = name
    tmpl.description = payload.get("description")
    tmpl.definition = payload.get("definition") or tmpl.definition
    db.add(tmpl)
    db.commit()
    _audit(db, action="UPDATE", resource_type="REPORT_TEMPLATE", resource_id=tmpl.id, details={"name": tmpl.name}, request=request)
    return JSONResponse({"ok": True})


@app.get("/api/dashboard")
def api_dashboard(
    template_id: int,
    db: Session = Depends(get_db),
    from_date: Optional[str] = Query(None),
    to_date: Optional[str] = Query(None),
    status: Optional[str] = Query("ALL"),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == template_id).first()
    tmpl = _ensure_template_has_widgets(db, tmpl)
    start_dt, end_dt = _parse_dates(from_date, to_date)
    metrics = _compute_metrics(db, start_dt, end_dt, status or "ALL")

    widgets = []
    for w in (tmpl.definition or {}).get("widgets", []):
        metric = w.get("metric")
        widgets.append({"type": w.get("type"), "metric": metric, "title": w.get("title"), "value": metrics.get(metric)})
    return JSONResponse({"template": {"id": tmpl.id, "name": tmpl.name}, "widgets": widgets})


@app.get("/api/customers/phone")
def lookup_customers_by_phone(
    phone: str,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    digits, _, _ = validate_lookup_inputs(phone=phone, cnic="", name="")
    if not digits:
        raise HTTPException(status_code=400, detail="Phone is required.")
    items = (
        db.query(Customer)
        .filter(Customer.is_deleted.is_(False))
        .filter(Customer.phone == digits)
        .order_by(Customer.name.asc())
        .limit(200)
        .all()
    )
    return JSONResponse({"count": len(items), "items": [_customer_to_dict(c) for c in items]})


@app.get("/api/customers/cnic")
def lookup_customers_by_cnic(
    cnic: str,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    _, digits, _ = validate_lookup_inputs(phone="", cnic=cnic, name="")
    if not digits:
        raise HTTPException(status_code=400, detail="CNIC is required.")
    normalized = format_cnic(digits)
    items = (
        db.query(Customer)
        .filter(Customer.is_deleted.is_(False))
        .filter((Customer.cnic == normalized) | (Customer.cnic == digits))
        .order_by(Customer.name.asc())
        .limit(200)
        .all()
    )
    return JSONResponse({"count": len(items), "items": [_customer_to_dict(c) for c in items]})


@app.get("/api/customers/name")
def lookup_customers_by_name(
    name: str,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    _, _, q = validate_lookup_inputs(phone="", cnic="", name=name)
    like = f"%{q}%"
    items = (
        db.query(Customer)
        .filter(Customer.is_deleted.is_(False))
        .filter((Customer.name.ilike(like)) | (Customer.business_name.ilike(like)))
        .order_by(Customer.name.asc())
        .limit(200)
        .all()
    )
    return JSONResponse({"count": len(items), "items": [_customer_to_dict(c) for c in items]})


@app.get("/api/customers/autocomplete")
def autocomplete_customer_names(
    query: str,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    q = (query or "").strip()
    if len(q) < 2:
        return JSONResponse({"items": []})
    like = f"%{q}%"
    rows = (
        db.query(Customer.name, Customer.business_name)
        .filter(Customer.is_deleted.is_(False))
        .filter((Customer.name.ilike(like)) | (Customer.business_name.ilike(like)))
        .order_by(Customer.name.asc())
        .limit(20)
        .all()
    )
    out: List[str] = []
    for name_value, business_name in rows:
        if name_value:
            value = str(name_value).strip()
            if value and value not in out:
                out.append(value)
        if business_name:
            value = str(business_name).strip()
            if value and value not in out:
                out.append(value)
    return JSONResponse({"items": out[:20]})


@app.get("/api/lookup/search")
def lookup_search(
    db: Session = Depends(get_db),
    phone: Optional[str] = Query(default=""),
    cnic: Optional[str] = Query(default=""),
    name: Optional[str] = Query(default=""),
    chassis: Optional[str] = Query(default=""),
    engine: Optional[str] = Query(default=""),
    sort_by: Optional[str] = Query(default="date"),
    sort_dir: Optional[str] = Query(default="desc"),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=20, ge=1, le=10000),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)

    try:
        phone_digits, cnic_digits, name_q, chassis_q, engine_q = validate_lookup_inputs(
            phone=phone or "", cnic=cnic or "", name=name or "", chassis=chassis or "", engine=engine or ""
        )
    except ValueError as exc:
        raise HTTPException(status_code=400, detail=str(exc))

    q = (
        db.query(
            Invoice.datetime.label("date_value"),
            Invoice.invoice_number.label("invoice_number"),
            Invoice.usin.label("usin"),
            Customer.name.label("customer_name"),
            Customer.father_name.label("father_name"),
            Customer.phone.label("mobile_number"),
            Customer.cnic.label("buyer_cnic"),
            Customer.ntn.label("buyer_ntn"),
            ProductModel.model_name.label("bike_model"),
            Motorcycle.chassis_number.label("chassis_number"),
            Motorcycle.engine_number.label("engine_number"),
            Invoice.sync_status.label("sync_status"),
            Invoice.fbr_invoice_number.label("fbr_invoice_number"),
            Invoice.total_sale_value.label("sale_value"),
            Invoice.total_quantity.label("quantity"),
            Invoice.status_updated_at.label("status_updated_at"),
        )
        .select_from(Invoice)
        .join(Customer, Invoice.customer_id == Customer.id, isouter=True)
        .join(InvoiceItem, InvoiceItem.invoice_id == Invoice.id, isouter=True)
        .join(Motorcycle, InvoiceItem.motorcycle_id == Motorcycle.id, isouter=True)
        .join(ProductModel, Motorcycle.product_model_id == ProductModel.id, isouter=True)
        .filter((Customer.is_deleted.is_(False)) | (Customer.id.is_(None)))
    )

    if phone_digits:
        q = q.filter(Customer.phone == phone_digits)
    if cnic_digits:
        normalized = format_cnic(cnic_digits)
        q = q.filter((Customer.cnic == normalized) | (Customer.cnic == cnic_digits))
    if name_q:
        like = f"%{name_q}%"
        q = q.filter((Customer.name.ilike(like)) | (Customer.business_name.ilike(like)))
    if chassis_q:
        like = f"%{chassis_q}%"
        q = q.filter(Motorcycle.chassis_number.ilike(like))
    if engine_q:
        like = f"%{engine_q}%"
        q = q.filter(Motorcycle.engine_number.ilike(like))

    sort_map = {
        "customer_name": Customer.name,
        "father_name": Customer.father_name,
        "mobile_number": Customer.phone,
        "invoice_number": Invoice.invoice_number,
        "bike_model": ProductModel.model_name,
        "chassis_number": Motorcycle.chassis_number,
        "engine_number": Motorcycle.engine_number,
        "date": Invoice.datetime,
    }
    sort_col = sort_map.get((sort_by or "date").strip(), Invoice.datetime)
    direction = (sort_dir or "desc").strip().lower()
    order_expr = asc(sort_col) if direction == "asc" else desc(sort_col)
    q = q.order_by(order_expr)

    total_count = int(q.order_by(None).with_entities(func.count()).scalar() or 0)
    total_pages = int((total_count + page_size - 1) / page_size) if total_count else 1
    offset = (page - 1) * page_size
    rows = q.offset(offset).limit(page_size).all()

    items: List[Dict[str, Any]] = []
    for r in rows:
        dt_value = getattr(r, "date_value", None)
        submission_value = getattr(r, "status_updated_at", None)
        raw_usin = getattr(r, "usin", None) or ""
        raw_fbr = getattr(r, "fbr_invoice_number", None) or ""
        effective_usin = raw_fbr or raw_usin or ""
        items.append(
            {
                "customer_name": getattr(r, "customer_name", None) or "",
                "father_name": getattr(r, "father_name", None) or "",
                "mobile_number": getattr(r, "mobile_number", None) or "",
                "buyer_cnic": getattr(r, "buyer_cnic", None) or "",
                "buyer_ntn": getattr(r, "buyer_ntn", None) or "",
                "invoice_number": getattr(r, "invoice_number", None) or "",
                "bike_model": getattr(r, "bike_model", None) or "",
                "chassis_number": getattr(r, "chassis_number", None) or "",
                "engine_number": getattr(r, "engine_number", None) or "",
                "date": dt_value.isoformat(sep=" ") if dt_value else "",
                "submission_date": submission_value.isoformat(sep=" ") if submission_value else "",
                "sync_status": getattr(r, "sync_status", None) or "",
                "fbr_invoice_number": raw_fbr,
                "usin": effective_usin,
                "sale_value": float(getattr(r, "sale_value", 0) or 0),
                "quantity": float(getattr(r, "quantity", 0) or 0),
            }
        )

    return JSONResponse(
        {
            "count": total_count,
            "page": page,
            "page_size": page_size,
            "total_pages": total_pages,
            "items": items,
        }
    )


@app.get("/api/invoices/{invoice_number}/details")
def api_invoice_details(
    invoice_number: str,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    inv_num = (invoice_number or "").strip()
    if not inv_num:
        raise HTTPException(status_code=400, detail="invoice_number required")

    inv = (
        db.query(Invoice)
        .options(
            joinedload(Invoice.customer),
            joinedload(Invoice.items).joinedload(InvoiceItem.motorcycle).joinedload(Motorcycle.product_model),
        )
        .filter(Invoice.invoice_number == inv_num)
        .first()
    )
    if not inv:
        raise HTTPException(status_code=404, detail="Invoice not found")
    return JSONResponse(invoice_to_detail_dict(inv))


@app.post("/api/invoices/{invoice_number}/retry")
def retry_invoice_upload(
    invoice_number: str,
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    inv_num = (invoice_number or "").strip()
    if not inv_num:
        raise HTTPException(status_code=400, detail="invoice_number required")

    with _retry_lock:
        now = time.time()
        last = _last_retry_at.get(inv_num, 0.0)
        if now - last < 30.0:
            raise HTTPException(status_code=429, detail="Too many retry attempts. Please wait and try again.")
        _last_retry_at[inv_num] = now

    inv = db.query(Invoice).filter(Invoice.invoice_number == inv_num).first()
    if not inv:
        raise HTTPException(status_code=404, detail="Invoice not found")
    if (inv.sync_status or "").upper() not in ["PENDING", "FAILED"]:
        return JSONResponse(
            {
                "invoice_number": inv.invoice_number,
                "sync_status": inv.sync_status,
                "message": "Retry not available for this status.",
            }
        )

    inv.fbr_response_message = "Manual retry requested."
    db.add(inv)
    db.commit()
    db.refresh(inv)

    try:
        invoice_service.sync_invoice(db, inv)
        db.commit()
        db.refresh(inv)
    except Exception as e:
        db.rollback()
        inv = db.query(Invoice).filter(Invoice.invoice_number == inv_num).first()
        status = inv.sync_status if inv else "PENDING"
        msg = str(e)
        return JSONResponse(
            {
                "invoice_number": inv_num,
                "sync_status": status,
                "message": msg,
            },
            status_code=200,
        )

    _audit(
        db,
        action="RETRY_SYNC",
        resource_type="INVOICE",
        resource_id=int(inv.id),
        details={"invoice_number": inv.invoice_number, "sync_status": inv.sync_status, "triggered_by": role},
        request=request,
    )

    return JSONResponse(
        {
            "invoice_number": inv.invoice_number,
            "sync_status": inv.sync_status,
            "message": inv.fbr_response_message or "",
        }
    )


@app.get("/export/{fmt}")
def export_report(
    fmt: str,
    template_id: int,
    request: Request,
    db: Session = Depends(get_db),
    from_date: Optional[str] = Query(None),
    to_date: Optional[str] = Query(None),
    status: Optional[str] = Query("ALL"),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> StreamingResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == template_id).first()
    tmpl = _ensure_template_has_widgets(db, tmpl)
    start_dt, end_dt = _parse_dates(from_date, to_date)
    metrics = _compute_metrics(db, start_dt, end_dt, status or "ALL")
    content, filename, media_type = _export_bytes(tmpl, metrics, fmt)
    _audit(db, action="EXPORT", resource_type="REPORT_TEMPLATE", resource_id=tmpl.id, details={"fmt": fmt, "filename": filename}, request=request)
    return StreamingResponse(BytesIO(content), media_type=media_type, headers={"Content-Disposition": f"attachment; filename={filename}"})


@app.get("/invoices.csv")
def invoices_csv_legacy(
    request: Request,
    db: Session = Depends(get_db),
    from_date: Optional[str] = Query(None),
    to_date: Optional[str] = Query(None),
    status: Optional[str] = Query("ALL"),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> StreamingResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)
    tmpl = _load_or_create_default_template(db)
    start_dt, end_dt = _parse_dates(from_date, to_date)
    metrics = _compute_metrics(db, start_dt, end_dt, status or "ALL")
    content = _export_csv(metrics)
    _audit(db, action="EXPORT", resource_type="REPORT_TEMPLATE", resource_id=tmpl.id, details={"fmt": "csv", "legacy": True}, request=request)
    return StreamingResponse(BytesIO(content), media_type="text/csv", headers={"Content-Disposition": "attachment; filename=invoices.csv"})


@app.get("/api/schedules")
def list_schedules(
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])
    schedules = db.query(ReportSchedule).order_by(ReportSchedule.id.desc()).all()
    items = []
    for s in schedules:
        tmpl = db.query(ReportTemplate).filter(ReportTemplate.id == s.template_id).first()
        items.append(
            {
                "id": s.id,
                "template_id": s.template_id,
                "template_name": tmpl.name if tmpl else "",
                "enabled": bool(s.enabled),
                "interval_minutes": int(s.interval_minutes or 60),
                "export_format": s.export_format,
                "recipients": s.recipients or [],
                "last_run_at": s.last_run_at.isoformat(sep=" ") if s.last_run_at else None,
            }
        )
    return JSONResponse({"items": items})


@app.post("/api/schedules")
def create_schedule(
    payload: Dict[str, Any],
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])
    template_id = int(payload.get("template_id") or 0)
    interval = int(payload.get("interval_minutes") or 60)
    export_format = (payload.get("export_format") or "pdf").lower()
    recipients = payload.get("recipients") or []
    schedule = ReportSchedule(
        template_id=template_id,
        interval_minutes=interval,
        export_format=export_format,
        recipients=recipients,
        enabled=bool(payload.get("enabled", True)),
        created_by_role=role,
    )
    db.add(schedule)
    db.commit()
    db.refresh(schedule)
    _audit(db, action="CREATE", resource_type="REPORT_SCHEDULE", resource_id=schedule.id, details={"template_id": template_id}, request=request)
    return JSONResponse({"id": schedule.id})


@app.put("/api/schedules/{schedule_id}")
def update_schedule(
    schedule_id: int,
    payload: Dict[str, Any],
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])
    sch = db.query(ReportSchedule).filter(ReportSchedule.id == schedule_id).first()
    if not sch:
        raise HTTPException(status_code=404, detail="Not found")
    if "enabled" in payload:
        sch.enabled = bool(payload.get("enabled"))
    if "interval_minutes" in payload:
        sch.interval_minutes = int(payload.get("interval_minutes") or sch.interval_minutes or 60)
    if "export_format" in payload:
        sch.export_format = (payload.get("export_format") or sch.export_format).lower()
    if "recipients" in payload:
        sch.recipients = payload.get("recipients")
    db.add(sch)
    db.commit()
    _audit(db, action="UPDATE", resource_type="REPORT_SCHEDULE", resource_id=sch.id, details={"enabled": sch.enabled}, request=request)
    return JSONResponse({"ok": True})


def _get_valid_finance_token(db: Session, token: str) -> Any:
    # FinancePortalToken is currently missing from models.py
    raise HTTPException(status_code=501, detail="Financing module is currently unavailable.")
    # t = (token or "").strip()
    # if not t:
    #     raise HTTPException(status_code=400, detail="Token is required.")
    # row = db.query(FinancePortalToken).filter(FinancePortalToken.token == t).first()
    # if not row or row.revoked_at is not None:
    #     raise HTTPException(status_code=404, detail="Invalid token.")
    # if row.expires_at and row.expires_at <= datetime.utcnow():
    #     raise HTTPException(status_code=403, detail="Token expired.")
    # return row


@app.get("/credit-portal", response_class=HTMLResponse)
def credit_portal_home() -> str:
    return """
    <!doctype html>
    <html lang="en">
      <head>
        <meta charset="utf-8"/>
        <meta name="viewport" content="width=device-width, initial-scale=1"/>
        <title>Credit Portal</title>
        <link href="https://cdn.jsdelivr.net/npm/bootstrap@5.3.3/dist/css/bootstrap.min.css" rel="stylesheet"/>
      </head>
      <body class="bg-light">
        <div class="container py-5">
          <div class="row justify-content-center">
            <div class="col-12 col-md-7 col-lg-5">
              <div class="card shadow-sm">
                <div class="card-body">
                  <h4 class="mb-2">Credit Portal</h4>
                  <p class="text-muted mb-4">Enter your access token to view your financing account.</p>
                  <form method="GET" action="/credit-portal/go">
                    <div class="mb-3">
                      <label class="form-label">Access Token</label>
                      <input class="form-control" name="token" placeholder="Paste token here" required/>
                    </div>
                    <button class="btn btn-primary w-100" type="submit">Continue</button>
                  </form>
                </div>
              </div>
              <div class="text-center text-muted small mt-3">
                Do not share your token with anyone.
              </div>
            </div>
          </div>
        </div>
      </body>
    </html>
    """


@app.get("/credit-portal/go")
def credit_portal_go(token: str = Query(default="")) -> RedirectResponse:
    t = (token or "").strip()
    if not t:
        return RedirectResponse("/credit-portal", status_code=302)
    return RedirectResponse(f"/credit-portal/{t}", status_code=302)


@app.get("/credit-portal/{token}", response_class=HTMLResponse)
def credit_portal_account(token: str, db: Session = Depends(get_db)) -> str:
    tok = _get_valid_finance_token(db, token)
    # The following code is unreachable due to the raise in _get_valid_finance_token
    return ""


@app.get("/credit-portal/{token}/loan/{loan_id}", response_class=HTMLResponse)
def credit_portal_loan(token: str, loan_id: int, db: Session = Depends(get_db)) -> str:
    tok = _get_valid_finance_token(db, token)
    # The following code is unreachable
    return ""


@app.get("/credit-portal/{token}/loan/{loan_id}/payment", response_class=HTMLResponse)
def credit_portal_payment(token: str, loan_id: int, db: Session = Depends(get_db)) -> str:
    tok = _get_valid_finance_token(db, token)
    # The following code is unreachable
    return ""


@app.post("/credit-portal/{token}/loan/{loan_id}/payment")
async def credit_portal_payment_submit(
    token: str,
    loan_id: int,
    request: Request,
    db: Session = Depends(get_db),
) -> RedirectResponse:
    tok = _get_valid_finance_token(db, token)
    # unreachable
    return RedirectResponse(f"/credit-portal/{token}", status_code=302)


# ---------------------------------------------------------------------------
# FastReport Integration Endpoints
# ---------------------------------------------------------------------------

@app.get("/api/fastreports/status")
def fastreports_status(
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    """Return whether FastReport Desktop is available and list its templates."""
    from app.services.fastreport_bridge import (
        find_fastreports,
        ensure_templates_dir,
        default_template_path,
    )

    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)

    info = find_fastreports()
    available = info is not None
    templates_dir = ensure_templates_dir()
    templates = []
    try:
        for entry in sorted(templates_dir.glob("*.frx")):
            templates.append({
                "name": entry.stem,
                "file_name": entry.name,
                "size_bytes": entry.stat().st_size,
                "modified_at": datetime.fromtimestamp(entry.stat().st_mtime).isoformat(sep=" "),
            })
    except Exception as exc:
        logger.exception("Failed to list FastReport templates: %s", exc)

    return JSONResponse({
        "available": available,
        "builder_exe": str(info.builder_exe) if info else None,
        "designer_exe": str(info.designer_exe) if info and info.designer_exe else None,
        "templates_dir": str(templates_dir),
        "templates": templates,
    })


@app.post("/api/fastreports/designer/open")
def fastreports_open_designer(
    payload: Dict[str, Any],
    request: Request,
    db: Session = Depends(get_db),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> JSONResponse:
    """Launch the FastReport Designer (FRDesigner.exe) for a .frx template."""
    from app.services.fastreport_bridge import open_designer, is_fastreports_available

    role = _get_role(x_user_role)
    _require_auth(x_api_key, role, required_roles=["admin", "manager"])

    if not is_fastreports_available():
        raise HTTPException(status_code=400, detail="FastReport Desktop is not installed.")

    template_name = (payload.get("template_name") or "").strip() or None
    ok, err = open_designer(template_name)
    if not ok:
        raise HTTPException(status_code=500, detail=err or "Failed to launch FastReport Designer.")

    _audit(
        db,
        action="OPEN_DESIGNER",
        resource_type="FASTREPORT_TEMPLATE",
        resource_id=0,
        details={"template_name": template_name or "<blank>"},
        request=request,
    )
    return JSONResponse({"ok": True})


@app.post("/api/fastreports/export/{template_name}")
def fastreports_direct_export(
    template_name: str,
    request: Request,
    db: Session = Depends(get_db),
    fmt: str = Query("pdf"),
    from_date: Optional[str] = Query(None),
    to_date: Optional[str] = Query(None),
    status: Optional[str] = Query("ALL"),
    x_api_key: Optional[str] = Header(default=None, alias="X-API-Key"),
    x_user_role: Optional[str] = Header(default=None, alias="X-User-Role"),
) -> StreamingResponse:
    """Direct export via FastReport using a named .frx template + live metrics."""
    from app.services.fastreport_bridge import (
        build_report,
        is_fastreports_available,
    )

    role = _get_role(x_user_role)
    _require_auth(x_api_key, role)

    if not is_fastreports_available():
        raise HTTPException(status_code=400, detail="FastReport Desktop is not installed.")

    start_dt, end_dt = _parse_dates(from_date, to_date)
    metrics = _compute_metrics(db, start_dt, end_dt, status or "ALL")

    report_data: Dict[str, Any] = {
        "report_title": template_name,
        "generated_at": datetime.utcnow().isoformat(sep=" "),
        "period_from": start_dt.isoformat(sep=" ") if start_dt else "",
        "period_to": end_dt.isoformat(sep=" ") if end_dt else "",
        "status_filter": status or "ALL",
        "total_invoices": metrics.get("total_invoices", 0),
        "total_amount": float(metrics.get("total_amount", 0.0)),
        "avg_invoice_amount": float(metrics.get("avg_invoice_amount", 0.0)),
        "daily_sales": metrics.get("daily_sales", []),
        "status_breakdown": metrics.get("status_breakdown", []),
        "invoices": metrics.get("invoices", []),
    }

    fmt_lower = (fmt or "pdf").lower()
    result = build_report(
        template_name_or_path=template_name,
        data=report_data,
        export_format=fmt_lower,
    )

    if not result.ok or not result.output_path:
        raise HTTPException(status_code=500, detail=result.error or "FastReport render failed.")

    out_path = result.output_path
    if not out_path.is_file():
        raise HTTPException(status_code=500, detail="FastReport did not produce output file.")

    suffix_map = {
        "pdf": ".pdf", "html": ".html", "htm": ".html",
        "xlsx": ".xlsx", "excel": ".xlsx", "csv": ".csv",
        "docx": ".docx", "rtf": ".rtf", "pptx": ".pptx",
        "png": ".png", "jpg": ".jpg",
    }
    media_type_map = {
        "pdf": "application/pdf", "html": "text/html", "htm": "text/html",
        "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "csv": "text/csv",
        "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "rtf": "application/rtf",
        "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "png": "image/png", "jpg": "image/jpeg",
    }
    suffix = suffix_map.get(fmt_lower, f".{fmt_lower}")
    media_type = media_type_map.get(fmt_lower, "application/octet-stream")
    filename = f"{template_name}_{datetime.utcnow().strftime('%Y%m%d_%H%M%S')}{suffix}"
    file_bytes = out_path.read_bytes()

    _audit(
        db,
        action="FR_EXPORT",
        resource_type="FASTREPORT_TEMPLATE",
        resource_id=0,
        details={"template_name": template_name, "fmt": fmt_lower, "filename": filename},
        request=request,
    )
    return StreamingResponse(
        BytesIO(file_bytes),
        media_type=media_type,
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )
